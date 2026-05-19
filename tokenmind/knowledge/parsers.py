"""Rich document parsers for the knowledge base.

Per-format extractors that preserve more structure than naive
``zipfile + itertext`` does — tables stay as ``cell | cell`` rows, PPTX
slides stay numbered, DOCX paragraphs keep their heading level. When a
``VLMConfig`` is provided, complex PDF pages and embedded Office images
get captioned by a vision-language model; without it the parsers
degrade cleanly to text-only extraction.

The legacy ``.doc`` / ``.ppt`` binary formats are converted to their
modern ZIP-based counterparts via the locally-installed LibreOffice
(``soffice``) — no Docker daemon required. Conversion fails fast with a
clear error when LibreOffice isn't on the system.
"""

from __future__ import annotations

import base64
import concurrent.futures
import io
import subprocess
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from loguru import logger

from tokenmind.utils.office import find_soffice


@dataclass
class VLMConfig:
    """Vision-language model client configuration for parsing-time image captioning."""

    model: str
    api_key: str
    api_base: str | None = None
    timeout: int = 30
    max_dim: int = 1280
    # Max number of concurrent VLM HTTP calls per document. Larger values
    # accelerate documents with many embedded images but push API spend higher.
    max_workers: int = 8


@dataclass
class ParsedPage:
    page_num: int
    content: str
    tokens_used: int = 0
    method: str = "text"


@dataclass
class ParsedDocument:
    file_name: str
    file_type: str
    pages: list[ParsedPage] = field(default_factory=list)

    @property
    def total_tokens(self) -> int:
        return sum(p.tokens_used for p in self.pages)

    def as_text(self) -> str:
        return "\n\n".join(p.content for p in self.pages if p.content.strip())


class LegacyOfficeConversionError(RuntimeError):
    """Raised when ``.doc`` / ``.ppt`` cannot be converted to the modern format."""


# Inline DOCX images get a low threshold because diagrams / sparklines that
# carry real information are often around 1–5KB. Source project value: 1KB.
_MIN_DOCX_IMAGE_BYTES = 1 * 1024  # 1KB
# PPTX images get a stricter cutoff — slides contain many decorative icons
# (bullet glyphs, watermark logos) that would burn VLM tokens for no gain.
_MIN_PPTX_IMAGE_BYTES = 5 * 1024  # 5KB
# Approximate per-segment character cap for DOCX paging.
_DOCX_SEGMENT_CHAR_LIMIT = 1000


# ---------------------------------------------------------------------------
# VLM image captioning
# ---------------------------------------------------------------------------

def _caption_image(
    image_bytes: bytes,
    vlm: VLMConfig,
    *,
    context_hint: str = "",
    strict_filtering: bool = True,
) -> tuple[str | None, int]:
    """Send a single image to the configured VLM, return ``(caption, tokens)``.

    ``strict_filtering=True`` (PPT mode) tells the model to drop decorative
    images outright by returning a sentinel that we map to ``None``. The
    non-strict mode (DOCX) keeps weak captions to maximise recall on inline
    diagrams.
    """
    try:
        from PIL import Image
        from openai import OpenAI
    except ImportError as exc:
        logger.warning("VLM dependencies missing: {}", exc)
        return None, 0

    try:
        img = Image.open(io.BytesIO(image_bytes))
        img.thumbnail((vlm.max_dim, vlm.max_dim))
        buf = io.BytesIO()
        save_format = "PNG" if img.mode in {"RGBA", "P"} else "JPEG"
        img.convert("RGB" if save_format == "JPEG" else img.mode).save(buf, format=save_format)
        data_url = (
            f"data:image/{save_format.lower()};base64,"
            f"{base64.b64encode(buf.getvalue()).decode('ascii')}"
        )
    except Exception:
        logger.exception("Failed to preprocess image for VLM")
        return None, 0

    if strict_filtering:
        instructions = (
            "你是文档解析助手。请用中文描述图片中的有用信息（图表数据、文字、"
            "示意图含义）。如果图片是纯装饰、logo、模糊背景或无可识别内容，"
            "只回复 SKIP。"
        )
    else:
        instructions = (
            "你是文档解析助手。请用中文简洁描述图片中的内容、表格或图表数据，"
            "重点提取文字与数值信息。"
        )

    user_prompt = instructions
    if context_hint:
        user_prompt = f"{instructions}\n\n上下文：{context_hint}"

    try:
        client = OpenAI(
            api_key=vlm.api_key,
            base_url=vlm.api_base or None,
            timeout=vlm.timeout,
        )
        resp = client.chat.completions.create(
            model=vlm.model,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": user_prompt},
                        {"type": "image_url", "image_url": {"url": data_url}},
                    ],
                }
            ],
            max_tokens=800,
        )
        text = (resp.choices[0].message.content or "").strip()
        tokens = getattr(resp.usage, "total_tokens", 0) if resp.usage else 0
        if strict_filtering and text.upper() == "SKIP":
            return None, tokens
        if not text:
            return None, tokens
        return text, tokens
    except Exception as exc:
        logger.warning("VLM call failed: {}", exc)
        return None, 0


def _caption_images_batch(
    tasks: list[tuple[str, bytes, str, bool]],
    vlm: VLMConfig,
) -> dict[str, tuple[str | None, int]]:
    """Fan ``tasks`` out across a thread pool and gather results by task_id.

    Each task is ``(task_id, image_bytes, context_hint, strict_filtering)``.
    The result dict maps task_id → ``(caption, tokens)`` exactly like the
    single-shot ``_caption_image`` helper, so callers can stay agnostic to
    whether captioning ran in parallel or not. Empty input short-circuits
    without spinning a thread pool.
    """
    if not tasks:
        return {}
    if vlm.max_workers <= 1 or len(tasks) == 1:
        return {
            task_id: _caption_image(
                blob, vlm, context_hint=ctx, strict_filtering=strict
            )
            for task_id, blob, ctx, strict in tasks
        }

    results: dict[str, tuple[str | None, int]] = {}
    with concurrent.futures.ThreadPoolExecutor(max_workers=vlm.max_workers) as executor:
        future_to_id = {
            executor.submit(
                _caption_image, blob, vlm, context_hint=ctx, strict_filtering=strict
            ): task_id
            for task_id, blob, ctx, strict in tasks
        }
        for future in concurrent.futures.as_completed(future_to_id):
            task_id = future_to_id[future]
            try:
                results[task_id] = future.result()
            except Exception as exc:
                logger.warning("VLM batch worker failed for {}: {}", task_id, exc)
                results[task_id] = (None, 0)
    return results


# ---------------------------------------------------------------------------
# Legacy format conversion via local LibreOffice
# ---------------------------------------------------------------------------

def _convert_legacy_to_modern(source: Path, target_ext: str) -> Path:
    """Convert ``.doc`` → ``.docx`` (or ``.ppt`` → ``.pptx``) using local
    LibreOffice. Returns the path to the converted file inside a tempdir.

    Raises ``LegacyOfficeConversionError`` if LibreOffice isn't installed or
    the conversion fails.
    """
    soffice = find_soffice()
    if not soffice:
        raise LegacyOfficeConversionError(
            "LibreOffice (soffice) not found. Install LibreOffice to enable "
            ".doc / .ppt parsing, or upload the document in the modern "
            "format (.docx / .pptx)."
        )

    target_format = target_ext.lstrip(".")  # 'docx' or 'pptx'
    out_dir = Path(tempfile.mkdtemp(prefix="tm-legacy-conv-"))
    try:
        result = subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                target_format,
                "--outdir",
                str(out_dir),
                str(source),
            ],
            capture_output=True,
            text=True,
            timeout=120,
        )
    except subprocess.TimeoutExpired as exc:
        raise LegacyOfficeConversionError(
            f"LibreOffice timed out converting {source.name}"
        ) from exc

    if result.returncode != 0:
        raise LegacyOfficeConversionError(
            f"LibreOffice failed to convert {source.name}: "
            f"{(result.stderr or result.stdout or '').strip()[:400]}"
        )

    converted = out_dir / f"{source.stem}.{target_format}"
    if not converted.exists():
        raise LegacyOfficeConversionError(
            f"LibreOffice produced no output for {source.name}"
        )
    return converted


# ---------------------------------------------------------------------------
# PDF (pymupdf / fitz backend)
# ---------------------------------------------------------------------------

# Fraction of the page area an image must cover to push the page into the
# VLM path. Mirrors smart-document-parser-API: large images = likely a
# chart / diagram that text extraction alone won't surface.
_PDF_IMAGE_AREA_RATIO = 0.05
# Below this many characters a page is considered "thin on text"; combined
# with the presence of any image, that's the cheap path to flag the page.
_PDF_TEXT_LIGHT_THRESHOLD = 800
# pymupdf pixmap render DPI when sending the page to a VLM. 100 dpi keeps
# images compact enough to stay under typical 1MB image-input limits while
# still letting the VLM read body-text-sized labels.
_PDF_VLM_RENDER_DPI = 100


def _pdf_page_is_complex(page: Any) -> bool:
    """fitz-backed complexity heuristic: trigger VLM when a page either has
    no extractable text but contains images, or has thin text and at least
    one image covering > 5% of the page area. Falls back to "any image with
    thin text" if bbox lookup fails (some embedded image refs raise inside
    fitz on damaged PDFs)."""
    text = (page.get_text() or "").strip()
    raw_images = page.get_images(full=True)
    if not raw_images:
        return False
    if not text:
        return True
    if len(text) >= _PDF_TEXT_LIGHT_THRESHOLD:
        return False
    page_rect = page.rect
    page_area = float(page_rect.width) * float(page_rect.height)
    if page_area <= 0:
        return True
    for img in raw_images:
        try:
            rect = page.get_image_bbox(img)
        except Exception:
            # Damaged xref → assume worth captioning rather than skipping
            return True
        area = float(rect.width) * float(rect.height)
        if page_area > 0 and (area / page_area) > _PDF_IMAGE_AREA_RATIO:
            return True
    return False


def _render_pdf_page_to_png(page: Any, dpi: int = _PDF_VLM_RENDER_DPI) -> bytes | None:
    """Render a fitz Page object to PNG bytes via the native pixmap API.
    Returns None on rendering failures so the caller can fall back to text."""
    try:
        pix = page.get_pixmap(dpi=dpi)
        return pix.tobytes("png")
    except Exception as exc:
        logger.warning("fitz failed to render PDF page: {}", exc)
        return None


def parse_pdf(path: Path, vlm: VLMConfig | None = None) -> ParsedDocument:
    """Parse a PDF using pymupdf (fitz). Text-only pages use the native
    extractor; image-dense / text-light pages get sent to the configured
    VLM in parallel via ``_caption_images_batch`` when VLM is enabled."""
    import fitz

    doc = fitz.open(str(path))
    try:
        result = ParsedDocument(file_name=path.name, file_type="pdf")

        # First pass: collect text per page and queue VLM tasks for the
        # complex ones. We keep the original text so the LLM still sees it
        # alongside the VLM caption in the final output.
        text_pages: list[tuple[int, str]] = []
        vlm_tasks: list[tuple[str, bytes, str, bool]] = []
        vlm_raw_text: dict[str, str] = {}

        for index, page in enumerate(doc):
            page_num = index + 1
            raw_text = (page.get_text() or "").strip()

            if vlm and _pdf_page_is_complex(page):
                png_bytes = _render_pdf_page_to_png(page)
                if png_bytes:
                    task_id = f"page_{page_num}"
                    vlm_tasks.append((task_id, png_bytes, raw_text[:200], False))
                    vlm_raw_text[task_id] = raw_text
                    continue

            if raw_text:
                text_pages.append((page_num, raw_text))

        captions: dict[str, tuple[str | None, int]] = {}
        if vlm_tasks and vlm:
            captions = _caption_images_batch(vlm_tasks, vlm)

        # Stitch back together in page order so chunkers preserve sequencing.
        rendered: dict[int, ParsedPage] = {}
        for page_num, raw_text in text_pages:
            rendered[page_num] = ParsedPage(
                page_num=page_num,
                content=f"--- Page {page_num} ---\n{raw_text}",
                method="text",
            )
        for task_id, (caption, tokens) in captions.items():
            page_num = int(task_id.split("_", 1)[1])
            raw_text = vlm_raw_text.get(task_id, "")
            if caption:
                body = caption
                if raw_text:
                    body = f"{caption}\n\n[原始文本]\n{raw_text}"
                rendered[page_num] = ParsedPage(
                    page_num=page_num,
                    content=f"--- Page {page_num} ---\n{body}",
                    tokens_used=tokens,
                    method="vlm",
                )
            elif raw_text:
                # VLM returned nothing useful — keep the text-only version
                # rather than dropping the page entirely.
                rendered[page_num] = ParsedPage(
                    page_num=page_num,
                    content=f"--- Page {page_num} ---\n{raw_text}",
                    method="vlm_failed",
                )

        for page_num in sorted(rendered.keys()):
            result.pages.append(rendered[page_num])
        return result
    finally:
        doc.close()


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def _extract_table_text(table: Any) -> str:
    """Render a DOCX table as ``cell | cell`` rows. Supports nested tables."""
    rows_text: list[str] = []
    for row in table.rows:
        cells_text: list[str] = []
        for cell in row.cells:
            paragraph_text = "\n".join(
                p.text.strip() for p in cell.paragraphs if p.text.strip()
            )
            for nested in getattr(cell, "tables", []) or []:
                nested_str = _extract_table_text(nested)
                if nested_str:
                    paragraph_text = (
                        f"{paragraph_text}\n{nested_str}" if paragraph_text else nested_str
                    )
            cells_text.append(paragraph_text)
        if any(cells_text):
            rows_text.append(" | ".join(cells_text))
    return "\n".join(rows_text)


def parse_docx(path: Path, vlm: VLMConfig | None = None) -> ParsedDocument:
    from docx import Document
    from docx.table import Table as DocxTable
    from docx.text.paragraph import Paragraph

    doc_obj = Document(str(path))
    result = ParsedDocument(file_name=path.name, file_type="docx")

    # Pre-index image relationship blobs so we can resolve inline embed IDs
    # without re-scanning the package each time.
    rid_to_blob: dict[str, bytes] = {}
    if vlm:
        for rel in doc_obj.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    rid_to_blob[rel.rId] = rel.target_part.blob
                except Exception:
                    continue

    structure: list[dict[str, Any]] = []
    pending_image_caches: dict[str, tuple[str, bytes]] = {}

    def _collect_inline_images(xml_str: str, ctx_hint: str) -> None:
        if not vlm or 'embed="' not in xml_str:
            return
        for part in xml_str.split('embed="')[1:]:
            rid = part.split('"', 1)[0]
            blob = rid_to_blob.get(rid)
            if not blob or len(blob) < _MIN_DOCX_IMAGE_BYTES:
                continue
            task_id = f"img_{len(pending_image_caches)}"
            pending_image_caches[task_id] = (ctx_hint, blob)
            structure.append({"type": "img", "id": task_id})

    for element in doc_obj.element.body:
        tag = element.tag.split("}")[-1]

        if tag == "p":
            para = Paragraph(element, doc_obj)
            text = para.text.strip()
            if text:
                style_name = (para.style.name or "").lower()
                if "heading" in style_name:
                    level_str = para.style.name.replace("Heading", "").strip()
                    level = int(level_str) if level_str.isdigit() else 2
                    structure.append({"type": "text", "val": f"{'#' * level} {text}"})
                else:
                    structure.append({"type": "text", "val": text})
            _collect_inline_images(element.xml, text[:50] if text else "Word document image")

        elif tag == "tbl":
            table = DocxTable(element, doc_obj)
            table_text = _extract_table_text(table)
            if table_text:
                structure.append({"type": "text", "val": table_text})
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        _collect_inline_images(
                            para._element.xml,
                            (para.text or "Word table image")[:50],
                        )

    captions: dict[str, tuple[str | None, int]] = {}
    if pending_image_caches and vlm:
        tasks = [
            (task_id, blob, ctx_hint, False)
            for task_id, (ctx_hint, blob) in pending_image_caches.items()
        ]
        captions = _caption_images_batch(tasks, vlm)

    chunk: list[str] = []
    chunk_chars = 0
    chunk_tokens = 0
    segment_num = 1

    def _flush() -> None:
        nonlocal chunk, chunk_chars, chunk_tokens, segment_num
        if not chunk:
            return
        body = "\n\n".join(chunk)
        result.pages.append(
            ParsedPage(
                page_num=segment_num,
                content=f"--- Segment {segment_num} ---\n{body}",
                tokens_used=chunk_tokens,
                method="docx",
            )
        )
        chunk = []
        chunk_chars = 0
        chunk_tokens = 0
        segment_num += 1

    for item in structure:
        if item["type"] == "text":
            piece = item["val"]
        else:
            caption, tokens = captions.get(item["id"], (None, 0))
            if not caption:
                continue
            piece = f"\n> **[图片内容]**:\n{caption}\n"
            chunk_tokens += tokens
        chunk.append(piece)
        chunk_chars += len(piece)
        if chunk_chars > _DOCX_SEGMENT_CHAR_LIMIT:
            _flush()

    _flush()
    return result


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def parse_pptx(path: Path, vlm: VLMConfig | None = None) -> ParsedDocument:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    result = ParsedDocument(file_name=path.name, file_type="pptx")

    # First pass: collect per-slide parts and queue every embedded image so
    # we can batch-caption them after walking all slides. Each part is either
    # a literal text string or an {"img": task_id} marker that gets resolved
    # in the second pass.
    slide_parts: list[list[Any]] = []
    image_tasks: list[tuple[str, bytes, str, bool]] = []

    for i, slide in enumerate(prs.slides):
        page_num = i + 1
        parts: list[Any] = []

        title_shape = slide.shapes.title if slide.shapes.title else None
        if title_shape and title_shape.text:
            parts.append(f"## {title_shape.text.strip()}")

        sorted_shapes = sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))
        for shape in sorted_shapes:
            if title_shape is not None and shape == title_shape:
                continue
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if vlm and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                except Exception:
                    blob = None
                if blob and len(blob) >= _MIN_PPTX_IMAGE_BYTES:
                    task_id = f"slide{page_num}_img{len(image_tasks)}"
                    ctx_hint = (
                        f"PPT Slide {page_num}: "
                        + " ".join(p for p in parts[-3:] if isinstance(p, str))
                    )
                    image_tasks.append((task_id, blob, ctx_hint, True))
                    parts.append({"img": task_id})
        slide_parts.append(parts)

    captions: dict[str, tuple[str | None, int]] = {}
    if image_tasks and vlm:
        captions = _caption_images_batch(image_tasks, vlm)

    # Second pass: stitch captions back into each slide in order.
    for i, parts in enumerate(slide_parts):
        page_num = i + 1
        rendered: list[str] = []
        page_tokens = 0
        for part in parts:
            if isinstance(part, str):
                rendered.append(part)
            else:
                caption, tokens = captions.get(part["img"], (None, 0))
                page_tokens += tokens
                if caption:
                    rendered.append(f"\n{caption}\n")
        body = "\n".join(rendered)
        result.pages.append(
            ParsedPage(
                page_num=page_num,
                content=f"--- Slide {page_num} ---\n{body}",
                tokens_used=page_tokens,
                method="pptx",
            )
        )

    return result


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------

def _format_xlsx_cell(value: Any) -> str:
    """Render a cell value as a tidy string for LLM consumption."""
    from datetime import date, datetime, time

    if value is None:
        return ""
    if isinstance(value, datetime):
        # Drop the trailing 00:00:00 for date-only cells so headers like
        # "2024-01-15" don't read as "2024-01-15 00:00:00".
        if value.hour == 0 and value.minute == 0 and value.second == 0:
            return value.strftime("%Y-%m-%d")
        return value.strftime("%Y-%m-%d %H:%M:%S")
    if isinstance(value, date):
        return value.strftime("%Y-%m-%d")
    if isinstance(value, time):
        return value.strftime("%H:%M:%S")
    if isinstance(value, float) and value.is_integer():
        # 12.0 -> "12" so integers stored as floats don't read ugly.
        return str(int(value))
    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"
    return str(value)


def parse_xlsx(path: Path) -> ParsedDocument:
    from openpyxl import load_workbook

    # read_only=False so we can see ws.merged_cells.ranges and forward-fill
    # merged header values into the rows below — the read_only iterator
    # otherwise returns ``None`` for every cell except the top-left of each
    # merge, which silently corrupts wide multi-row headers. For typical
    # knowledge-base spreadsheets (under tens of MB) the memory cost is
    # acceptable; very large workbooks degrade ungracefully but still parse.
    wb = load_workbook(filename=str(path), data_only=True, read_only=False)
    result = ParsedDocument(file_name=path.name, file_type="xlsx")
    try:
        for sheet_index, sheet_name in enumerate(wb.sheetnames):
            ws = wb[sheet_name]

            # Build a (row, col) -> top-left value map so we can fill in the
            # cells that openpyxl reports as None inside a merged range.
            merge_fill: dict[tuple[int, int], Any] = {}
            for merge_range in ws.merged_cells.ranges:
                anchor_value = ws.cell(merge_range.min_row, merge_range.min_col).value
                for r in range(merge_range.min_row, merge_range.max_row + 1):
                    for c in range(merge_range.min_col, merge_range.max_col + 1):
                        merge_fill[(r, c)] = anchor_value

            lines: list[str] = []
            for row_idx, row in enumerate(ws.iter_rows(values_only=False), start=1):
                cells: list[str] = []
                for col_idx, cell in enumerate(row, start=1):
                    val = cell.value
                    if val is None and (row_idx, col_idx) in merge_fill:
                        val = merge_fill[(row_idx, col_idx)]
                    cells.append(_format_xlsx_cell(val))
                # Trim trailing empty cells so a sparse row doesn't render as
                # "x | | | | | | |".
                while cells and not cells[-1].strip():
                    cells.pop()
                if cells:
                    lines.append(" | ".join(cells))
            if not lines:
                continue
            body = "\n".join(lines)
            result.pages.append(
                ParsedPage(
                    page_num=sheet_index + 1,
                    content=f"--- Sheet: {sheet_name} ---\n{body}",
                    method="xlsx",
                )
            )
    finally:
        wb.close()
    return result


# ---------------------------------------------------------------------------
# Top-level dispatch
# ---------------------------------------------------------------------------

# Extensions we have first-class structured parsers for. Everything else
# (txt, md, json, …) falls back to plain UTF-8 reading.
_RICH_SUFFIXES = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls"}


def can_parse(suffix: str) -> bool:
    return suffix.lower() in _RICH_SUFFIXES


def extract_document_text(path: Path, vlm: VLMConfig | None = None) -> str:
    """Top-level dispatch used by ``KnowledgeService``. Returns concatenated
    text suitable for chunking + embedding."""
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return parse_pdf(path, vlm).as_text()
    if suffix == ".docx":
        return parse_docx(path, vlm).as_text()
    if suffix == ".doc":
        converted = _convert_legacy_to_modern(path, ".docx")
        try:
            return parse_docx(converted, vlm).as_text()
        finally:
            _cleanup_legacy_conversion(converted)
    if suffix == ".pptx":
        return parse_pptx(path, vlm).as_text()
    if suffix == ".ppt":
        converted = _convert_legacy_to_modern(path, ".pptx")
        try:
            return parse_pptx(converted, vlm).as_text()
        finally:
            _cleanup_legacy_conversion(converted)
    if suffix in {".xlsx", ".xls"}:
        # ``.xls`` is the old OLE binary format that openpyxl can't read; do a
        # soffice conversion first.
        if suffix == ".xls":
            converted = _convert_legacy_to_modern(path, ".xlsx")
            try:
                return parse_xlsx(converted).as_text()
            finally:
                _cleanup_legacy_conversion(converted)
        return parse_xlsx(path).as_text()
    return path.read_text(encoding="utf-8", errors="ignore")


def _cleanup_legacy_conversion(converted: Path) -> None:
    try:
        converted.unlink(missing_ok=True)
        converted.parent.rmdir()
    except OSError:
        pass
