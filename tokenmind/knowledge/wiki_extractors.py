"""Text extraction from various source formats."""
from __future__ import annotations

from pathlib import Path

TEXT_SUFFIXES = {".md", ".txt", ".markdown", ".rst", ".log"}


def extract_text(path: Path) -> str:
    """Dispatch by suffix; return UTF-8 text or raise."""
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix == ".docx":
        return _extract_docx(path)
    if suffix == ".pptx":
        return _extract_pptx(path)
    if suffix == ".xlsx":
        return _extract_xlsx(path)
    if suffix in TEXT_SUFFIXES:
        return path.read_text(encoding="utf-8", errors="replace")
    raise ValueError(f"unsupported file type: {suffix}")


def _extract_pdf(path: Path) -> str:
    import pypdf
    reader = pypdf.PdfReader(str(path))
    return "\n\n".join(p.extract_text() or "" for p in reader.pages)


def _extract_docx(path: Path) -> str:
    try:
        import docx  # python-docx
    except ImportError as exc:
        raise RuntimeError("python-docx not installed") from exc
    d = docx.Document(str(path))
    return "\n\n".join(p.text for p in d.paragraphs if p.text)


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx not installed") from exc
    pres = Presentation(str(path))
    out = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text:
                        out.append(text)
    return "\n".join(out)


def _extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
    except ImportError as exc:
        raise RuntimeError("openpyxl not installed") from exc
    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    rows = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            rows.append("\t".join(str(c) if c is not None else "" for c in row))
    return "\n".join(rows)
