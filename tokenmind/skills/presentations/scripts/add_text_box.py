#!/usr/bin/env python3
"""Add a free-form text box to a slide.

Use this for slides that don't fit a layout placeholder shape — callouts,
chart annotations, labels, body paragraphs on a Blank-layout slide.

Coordinates and sizes are in inches. The origin is the top-left of the
slide; positive x goes right, positive y goes down.

Examples
--------

  # Big heading on slide 1, top-left
  python add_text_box.py /tmp/deck.pptx --slide 0 \\
      --left 0.5 --top 0.4 --width 9 --height 1.0 \\
      --text "Quarterly Results — Q3" --font-size 36 --bold

  # Body paragraph below, gray, justified left
  python add_text_box.py /tmp/deck.pptx --slide 0 \\
      --left 0.5 --top 1.6 --width 9 --height 4 \\
      --text "We achieved record revenue this quarter…" \\
      --font-size 18 --font-color 4B5563

  # Multi-paragraph (newline-separated)
  python add_text_box.py /tmp/deck.pptx --slide 1 \\
      --left 0.5 --top 1.0 --width 9 --height 4 \\
      --text $'Revenue: $12.4M\\nGrowth: +18% YoY\\nChurn: 2.1%' \\
      --font-size 22

  # From a file (preserves all newlines verbatim)
  python add_text_box.py /tmp/deck.pptx --slide 2 \\
      --left 0.5 --top 1.0 --width 12 --height 5 \\
      --text-file body.txt --font-size 16
"""
from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:  # pragma: no cover - guarded by skill requires.python
    print(
        "Error: python-pptx is not installed. Run: pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(2)


_HEX_RE = re.compile(r"^[0-9A-Fa-f]{6}$")
_ALIGN = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


def _hex_to_rgb(hex_str: str) -> RGBColor:
    v = hex_str.lstrip("#")
    if not _HEX_RE.match(v):
        raise ValueError(f"--font-color must be 6-digit hex, got {hex_str!r}")
    return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def add_text_box(
    pptx: Path,
    *,
    slide_index: int,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    text: str,
    font_size: int | None = None,
    font_color: str | None = None,
    bold: bool = False,
    italic: bool = False,
    align: str | None = None,
) -> None:
    """Place a text box on ``slide_index``. Mutates the file in place."""
    prs = Presentation(pptx)
    if not 0 <= slide_index < len(prs.slides):
        raise ValueError(
            f"--slide must be in range 0..{len(prs.slides) - 1}, got {slide_index}"
        )

    color = _hex_to_rgb(font_color) if font_color else None

    if align is not None and align not in _ALIGN:
        raise ValueError(f"--align must be one of {sorted(_ALIGN)}, got {align!r}")

    slide = prs.slides[slide_index]
    box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = box.text_frame
    tf.word_wrap = True

    lines = text.split("\n")
    # First line uses the default paragraph python-pptx already created.
    first = tf.paragraphs[0]
    first.text = lines[0]
    if align:
        first.alignment = _ALIGN[align]
    if first.runs:
        run = first.runs[0]
        if font_size is not None:
            run.font.size = Pt(font_size)
        if bold:
            run.font.bold = True
        if italic:
            run.font.italic = True
        if color is not None:
            run.font.color.rgb = color

    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        if align:
            p.alignment = _ALIGN[align]
        if p.runs:
            run = p.runs[0]
            if font_size is not None:
                run.font.size = Pt(font_size)
            if bold:
                run.font.bold = True
            if italic:
                run.font.italic = True
            if color is not None:
                run.font.color.rgb = color

    prs.save(pptx)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("pptx", type=Path, help="Path to the .pptx file")
    parser.add_argument("--slide", type=int, required=True, help="Zero-based slide index")
    parser.add_argument("--left", type=float, required=True, help="Left position (inches)")
    parser.add_argument("--top", type=float, required=True, help="Top position (inches)")
    parser.add_argument("--width", type=float, required=True, help="Box width (inches)")
    parser.add_argument("--height", type=float, required=True, help="Box height (inches)")

    text_src = parser.add_mutually_exclusive_group(required=True)
    text_src.add_argument("--text", help="Text content. \\n for line breaks.")
    text_src.add_argument("--text-file", type=Path, help="Read text from a file (preserves newlines)")

    parser.add_argument("--font-size", type=int, default=None, help="Font size in points")
    parser.add_argument("--font-color", default=None, help="6-digit hex like 4B5563 (no '#')")
    parser.add_argument("--bold", action="store_true")
    parser.add_argument("--italic", action="store_true")
    parser.add_argument("--align", choices=sorted(_ALIGN), default=None)
    args = parser.parse_args(argv)

    if not args.pptx.is_file():
        print(f"Error: pptx not found: {args.pptx}", file=sys.stderr)
        return 1

    if args.text is not None:
        text = args.text
    else:
        if not args.text_file.is_file():
            print(f"Error: --text-file not found: {args.text_file}", file=sys.stderr)
            return 1
        text = args.text_file.read_text(encoding="utf-8")

    try:
        add_text_box(
            args.pptx,
            slide_index=args.slide,
            left_in=args.left,
            top_in=args.top,
            width_in=args.width,
            height_in=args.height,
            text=text,
            font_size=args.font_size,
            font_color=args.font_color,
            bold=args.bold,
            italic=args.italic,
            align=args.align,
        )
    except ValueError as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1

    print(f"Added text box on slide {args.slide} of {args.pptx}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
