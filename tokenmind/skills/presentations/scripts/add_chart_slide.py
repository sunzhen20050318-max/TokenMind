#!/usr/bin/env python3
"""Add a native PowerPoint chart to an existing slide.

The chart is editable inside PowerPoint (right-click → Edit Data). It is
*not* an embedded image — it's a real OOXML chart object with its own
data, legend, and series.

Examples
--------

  # Bar chart with two series
  python add_chart_slide.py /tmp/deck.pptx --slide 1 \\
      --left 0.5 --top 1.5 --width 9 --height 5 \\
      --type bar \\
      --data '{"categories": ["Q1","Q2","Q3","Q4"], \\
               "series": {"Revenue":[120,140,170,200], \\
                          "Costs":[90,100,110,125]}}' \\
      --title "FY24 Quarterly"

  # Line chart, single series
  python add_chart_slide.py /tmp/deck.pptx --slide 2 \\
      --left 1 --top 1 --width 11 --height 5 \\
      --type line \\
      --data '{"categories":["2020","2021","2022","2023","2024"], \\
               "series":{"ARR":[12,18,27,38,52]}}' \\
      --title "ARR Growth ($M)"

  # Pie chart (single series only)
  python add_chart_slide.py /tmp/deck.pptx --slide 3 \\
      --left 2 --top 1 --width 8 --height 5 \\
      --type pie \\
      --data '{"categories":["Direct","Partner","Self-serve"], \\
               "series":{"Mix":[55,30,15]}}'

Data spec
---------

  Either keys are mandatory:

  - ``categories`` — list of category labels (X axis / pie segment names)
  - ``series``     — dict mapping series name → list of values
                     (same length as ``categories``)

  Pie charts use the FIRST series only (extra series are silently dropped).
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches
except ImportError:  # pragma: no cover - guarded by skill requires.python
    print(
        "Error: python-pptx is not installed. Run: pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(2)


_CHART_TYPES = {
    "bar":    XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line":   XL_CHART_TYPE.LINE,
    "pie":    XL_CHART_TYPE.PIE,
}


def add_chart_slide(
    pptx: Path,
    *,
    slide_index: int,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    chart_type: str,
    data: dict,
    title: str | None = None,
) -> None:
    """Insert a chart on ``slide_index``."""
    if chart_type not in _CHART_TYPES:
        raise ValueError(
            f"--type must be one of {sorted(_CHART_TYPES)}, got {chart_type!r}"
        )
    if not isinstance(data, dict):
        raise ValueError("--data must decode to a JSON object")
    if "categories" not in data or "series" not in data:
        raise ValueError("--data must contain 'categories' and 'series' keys")
    if not isinstance(data["categories"], list):
        raise ValueError("--data.categories must be a list")
    if not isinstance(data["series"], dict) or not data["series"]:
        raise ValueError("--data.series must be a non-empty object")

    n_cats = len(data["categories"])
    for name, values in data["series"].items():
        if not isinstance(values, list):
            raise ValueError(f"Series {name!r} values must be a list")
        if len(values) != n_cats:
            raise ValueError(
                f"Series {name!r}: {len(values)} values does not match "
                f"{n_cats} categories"
            )

    prs = Presentation(pptx)
    if not 0 <= slide_index < len(prs.slides):
        raise ValueError(
            f"--slide must be in range 0..{len(prs.slides) - 1}, got {slide_index}"
        )
    slide = prs.slides[slide_index]

    chart_data = CategoryChartData()
    chart_data.categories = list(data["categories"])

    series_items = list(data["series"].items())
    if chart_type == "pie":
        # Pie: just the first series
        name, values = series_items[0]
        chart_data.add_series(name, values)
    else:
        for name, values in series_items:
            chart_data.add_series(name, values)

    graphic_frame = slide.shapes.add_chart(
        _CHART_TYPES[chart_type],
        Inches(left_in),
        Inches(top_in),
        Inches(width_in),
        Inches(height_in),
        chart_data,
    )

    if title:
        chart = graphic_frame.chart
        chart.has_title = True
        chart.chart_title.text_frame.text = title

    prs.save(pptx)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("pptx", type=Path, help="Path to the .pptx file")
    parser.add_argument("--slide", type=int, required=True, help="Zero-based slide index")
    parser.add_argument("--left", type=float, required=True, help="Left position (inches)")
    parser.add_argument("--top", type=float, required=True, help="Top position (inches)")
    parser.add_argument("--width", type=float, required=True, help="Chart width (inches)")
    parser.add_argument("--height", type=float, required=True, help="Chart height (inches)")
    parser.add_argument("--type", required=True, choices=sorted(_CHART_TYPES), help="Chart type")
    parser.add_argument(
        "--data",
        required=True,
        help='JSON: {"categories":[...], "series":{"Name":[v1,v2,...]}}',
    )
    parser.add_argument("--title", default=None, help="Chart title")
    args = parser.parse_args(argv)

    if not args.pptx.is_file():
        print(f"Error: pptx not found: {args.pptx}", file=sys.stderr)
        return 1

    try:
        data = json.loads(args.data)
    except json.JSONDecodeError as e:
        print(f"Error: --data is not valid JSON: {e}", file=sys.stderr)
        return 1

    try:
        add_chart_slide(
            args.pptx,
            slide_index=args.slide,
            left_in=args.left,
            top_in=args.top,
            width_in=args.width,
            height_in=args.height,
            chart_type=args.type,
            data=data,
            title=args.title,
        )
    except ValueError as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1

    print(f"Added {args.type} chart on slide {args.slide} of {args.pptx}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
