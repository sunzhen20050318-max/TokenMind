#!/usr/bin/env python3
"""Create a new .pptx presentation.

By default produces a 16:9 deck (13.333 in × 7.5 in). Use ``--size 4:3`` for
the classic 10 × 7.5 in shape, or ``--width-in`` / ``--height-in`` for a
fully custom size.

A new deck starts empty (no slides). Use ``add_slide.py`` to populate it.
``--cover-title`` is a convenience: if set, the first slide will be a
Title Layout with that text — useful for "give me a starting deck".

Examples
--------

  # Empty 16:9 deck
  python build_pptx.py --out /tmp/deck.pptx

  # Empty 4:3 deck
  python build_pptx.py --out /tmp/deck.pptx --size 4:3

  # Custom A4 landscape (in inches)
  python build_pptx.py --out /tmp/deck.pptx --width-in 11.69 --height-in 8.27

  # Quick-start with a title slide
  python build_pptx.py --out /tmp/deck.pptx --cover-title "Q3 Results"
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:  # pragma: no cover - guarded by skill requires.python
    print(
        "Error: python-pptx is not installed. Run: pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(2)


_SIZE_PRESETS = {
    "16:9": (13.333, 7.5),
    "4:3":  (10.0,  7.5),
}


def build_pptx(
    out: Path,
    *,
    size: str = "16:9",
    width_in: float | None = None,
    height_in: float | None = None,
    cover_title: str | None = None,
    overwrite: bool = False,
) -> Path:
    """Create a new .pptx. Returns the output path."""
    if out.exists() and not overwrite:
        raise FileExistsError(
            f"{out} already exists. Pass --overwrite to replace it."
        )
    out.parent.mkdir(parents=True, exist_ok=True)

    if width_in is not None or height_in is not None:
        if width_in is None or height_in is None:
            raise ValueError("Provide both --width-in and --height-in, or use --size")
        w_in, h_in = width_in, height_in
    elif size in _SIZE_PRESETS:
        w_in, h_in = _SIZE_PRESETS[size]
    else:
        raise ValueError(
            f"--size must be one of {sorted(_SIZE_PRESETS)} or use --width-in/--height-in, got {size!r}"
        )

    prs = Presentation()
    prs.slide_width = Inches(w_in)
    prs.slide_height = Inches(h_in)

    # Drop any pre-existing slides from the default template so callers
    # get a clean deck (python-pptx ships with zero slides anyway, but be
    # defensive in case a user supplies a template later).
    # Slides removal is fiddly in python-pptx; just skip if none exist.

    if cover_title:
        # Layout 0 is "Title Slide" in the default template
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = cover_title

    prs.save(out)
    return out


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("--out", required=True, type=Path, help="Path to write the .pptx file")
    parser.add_argument(
        "--size",
        default="16:9",
        choices=sorted(_SIZE_PRESETS),
        help="Slide aspect ratio (default: 16:9)",
    )
    parser.add_argument("--width-in", type=float, default=None, help="Custom slide width in inches")
    parser.add_argument("--height-in", type=float, default=None, help="Custom slide height in inches")
    parser.add_argument("--cover-title", default=None, help="Optional title-slide text")
    parser.add_argument("--overwrite", action="store_true", help="Overwrite the output file if it exists")
    args = parser.parse_args(argv)

    try:
        path = build_pptx(
            args.out,
            size=args.size,
            width_in=args.width_in,
            height_in=args.height_in,
            cover_title=args.cover_title,
            overwrite=args.overwrite,
        )
    except (FileExistsError, ValueError) as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1
    print(f"Created {path}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
