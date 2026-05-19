#!/usr/bin/env python3
"""Restyle an existing deck against a design-system spec.

The spec is a JSON object describing typography + colors. We walk every
slide and rewrite the run-level font/color of titles, body placeholders,
table cells, and free-form text boxes — *without* recreating the deck or
losing layout/positioning.

Examples
--------

  # Apply an inline design system
  python apply_design_system.py /tmp/deck.pptx --design '{
    "heading_font": "Inter",
    "body_font":    "Inter",
    "heading_color": "111827",
    "body_color":    "4B5563",
    "accent_color":  "2563EB",
    "slide_bg":      "FFFFFF"
  }'

  # Or load from a file (preferred for anything non-trivial)
  python apply_design_system.py /tmp/deck.pptx --design-file ./brand.json

Spec keys
---------

  All keys are OPTIONAL. Missing keys leave the corresponding aspect alone.

  - heading_font   : font family applied to slide titles + table header cells
  - body_font      : font family applied to everything else
  - heading_color  : 6-digit hex (no '#'), titles & table headers
  - body_color     : 6-digit hex, body / placeholder / textbox text
  - accent_color   : 6-digit hex, applied to chart series fills if any
  - slide_bg       : 6-digit hex, fills every slide's background

The script never DELETES or REORDERS shapes — it only restyles text runs
and adds background fills. Safe to re-run.
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
except ImportError:  # pragma: no cover - guarded by skill requires.python
    print(
        "Error: python-pptx is not installed. Run: pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(2)


_HEX_RE = re.compile(r"^[0-9A-Fa-f]{6}$")
_VALID_KEYS = {
    "heading_font",
    "body_font",
    "heading_color",
    "body_color",
    "accent_color",
    "slide_bg",
}


def _hex_to_rgb(value: str, label: str) -> RGBColor:
    v = value.lstrip("#")
    if not _HEX_RE.match(v):
        raise ValueError(f"{label} must be 6-digit hex, got {value!r}")
    return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def _restyle_text_frame(tf, *, font: str | None, color: RGBColor | None) -> int:
    """Apply font/color to every run in ``tf``. Returns number of runs touched."""
    touched = 0
    for para in tf.paragraphs:
        for run in para.runs:
            if font is not None:
                run.font.name = font
            if color is not None:
                run.font.color.rgb = color
            touched += 1
    return touched


def _is_title_shape(shape, slide) -> bool:
    return shape == slide.shapes.title


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def apply_design_system(pptx: Path, spec: dict[str, Any]) -> dict[str, int]:
    """Apply a design-system spec to ``pptx`` in place. Returns counters."""
    if not isinstance(spec, dict):
        raise ValueError("--design must decode to a JSON object")

    unknown = set(spec.keys()) - _VALID_KEYS
    if unknown:
        raise ValueError(
            f"Unknown design-system keys: {sorted(unknown)}. "
            f"Valid keys: {sorted(_VALID_KEYS)}"
        )

    heading_font = spec.get("heading_font")
    body_font = spec.get("body_font")
    heading_color = _hex_to_rgb(spec["heading_color"], "heading_color") if spec.get("heading_color") else None
    body_color = _hex_to_rgb(spec["body_color"], "body_color") if spec.get("body_color") else None
    accent_color = _hex_to_rgb(spec["accent_color"], "accent_color") if spec.get("accent_color") else None
    slide_bg = _hex_to_rgb(spec["slide_bg"], "slide_bg") if spec.get("slide_bg") else None

    prs = Presentation(pptx)
    counts = {"slides": 0, "titles": 0, "bodies": 0, "table_cells": 0, "charts": 0, "backgrounds": 0}

    for slide in prs.slides:
        counts["slides"] += 1
        if slide_bg is not None:
            _set_slide_bg(slide, slide_bg)
            counts["backgrounds"] += 1

        for shape in slide.shapes:
            if shape.has_text_frame:
                is_title = _is_title_shape(shape, slide)
                font = heading_font if is_title else body_font
                color = heading_color if is_title else body_color
                touched = _restyle_text_frame(shape.text_frame, font=font, color=color)
                if touched:
                    if is_title:
                        counts["titles"] += 1
                    else:
                        counts["bodies"] += 1

            elif shape.has_table:
                table = shape.table
                for r_idx, row in enumerate(table.rows):
                    for cell in row.cells:
                        is_header = (r_idx == 0)
                        font = heading_font if is_header else body_font
                        color = heading_color if is_header else body_color
                        _restyle_text_frame(cell.text_frame, font=font, color=color)
                        counts["table_cells"] += 1

            elif shape.has_chart and accent_color is not None:
                # Accent color hint: paint the first series with the accent fill.
                # Charts are complex — be conservative and only touch series 0.
                chart = shape.chart
                if chart.series:
                    series = chart.series[0]
                    fill = series.format.fill
                    fill.solid()
                    fill.fore_color.rgb = accent_color
                    counts["charts"] += 1

    prs.save(pptx)
    return counts


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("pptx", type=Path, help="Path to the .pptx file")

    src = parser.add_mutually_exclusive_group(required=True)
    src.add_argument("--design", help="Inline JSON spec")
    src.add_argument("--design-file", type=Path, help="Path to a JSON spec file")

    args = parser.parse_args(argv)

    if not args.pptx.is_file():
        print(f"Error: pptx not found: {args.pptx}", file=sys.stderr)
        return 1

    try:
        if args.design:
            spec = json.loads(args.design)
        else:
            if not args.design_file.is_file():
                print(f"Error: --design-file not found: {args.design_file}", file=sys.stderr)
                return 1
            spec = json.loads(args.design_file.read_text(encoding="utf-8"))
    except json.JSONDecodeError as e:
        print(f"Error: design spec is not valid JSON: {e}", file=sys.stderr)
        return 1

    try:
        counts = apply_design_system(args.pptx, spec)
    except ValueError as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1

    print(f"Applied design system to {args.pptx}: {counts}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
