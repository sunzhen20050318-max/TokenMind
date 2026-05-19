#!/usr/bin/env python3
"""Add an image to an existing slide.

Supports PNG, JPEG, GIF, BMP, TIFF, WMF, EMF — anything python-pptx's
``add_picture`` accepts. Hosting code in this skill stays in PNG/JPEG
territory because that's what visual-verification renderers handle.

Sizing rules:
  - If both ``--width`` and ``--height`` are given, the image is forced to
    that exact box (may distort).
  - If only one dimension is given, the other is computed to preserve the
    original aspect ratio.
  - If neither is given, the image uses its intrinsic EMU size (no scaling).

Examples
--------

  # Place a logo, preserve aspect, height = 1 inch
  python add_image_slide.py /tmp/deck.pptx --slide 0 \\
      --image ./logo.png --left 0.4 --top 0.4 --height 1.0

  # Full-bleed hero image (force fit)
  python add_image_slide.py /tmp/deck.pptx --slide 1 \\
      --image ./hero.jpg --left 0 --top 0 --width 13.333 --height 7.5
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:  # pragma: no cover - guarded by skill requires.python
    print(
        "Error: python-pptx is not installed. Run: pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(2)


def add_image_slide(
    pptx: Path,
    *,
    slide_index: int,
    image: Path,
    left_in: float,
    top_in: float,
    width_in: float | None = None,
    height_in: float | None = None,
) -> None:
    """Place ``image`` on ``slide_index``."""
    if not image.is_file():
        raise FileNotFoundError(image)

    prs = Presentation(pptx)
    if not 0 <= slide_index < len(prs.slides):
        raise ValueError(
            f"--slide must be in range 0..{len(prs.slides) - 1}, got {slide_index}"
        )
    slide = prs.slides[slide_index]

    kwargs: dict = {}
    if width_in is not None:
        kwargs["width"] = Inches(width_in)
    if height_in is not None:
        kwargs["height"] = Inches(height_in)

    slide.shapes.add_picture(
        str(image),
        Inches(left_in),
        Inches(top_in),
        **kwargs,
    )
    prs.save(pptx)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("pptx", type=Path, help="Path to the .pptx file")
    parser.add_argument("--slide", type=int, required=True, help="Zero-based slide index")
    parser.add_argument("--image", type=Path, required=True, help="Path to the image file")
    parser.add_argument("--left", type=float, required=True, help="Left position (inches)")
    parser.add_argument("--top", type=float, required=True, help="Top position (inches)")
    parser.add_argument("--width", type=float, default=None, help="Width (inches). Omit to preserve aspect.")
    parser.add_argument("--height", type=float, default=None, help="Height (inches). Omit to preserve aspect.")
    args = parser.parse_args(argv)

    if not args.pptx.is_file():
        print(f"Error: pptx not found: {args.pptx}", file=sys.stderr)
        return 1

    try:
        add_image_slide(
            args.pptx,
            slide_index=args.slide,
            image=args.image,
            left_in=args.left,
            top_in=args.top,
            width_in=args.width,
            height_in=args.height,
        )
    except (FileNotFoundError, ValueError) as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1

    print(f"Added image {args.image} on slide {args.slide} of {args.pptx}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
