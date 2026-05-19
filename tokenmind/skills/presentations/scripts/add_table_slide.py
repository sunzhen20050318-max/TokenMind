#!/usr/bin/env python3
"""Add a table to an existing slide.

Tables are native PowerPoint tables (editable in PowerPoint), not images.
The script lays out a uniformly-sized grid and fills it with the strings
or numbers you pass in via ``--data``.

Examples
--------

  # 3x2 metrics table
  python add_table_slide.py /tmp/deck.pptx --slide 1 \\
      --left 0.5 --top 2 --width 9 --height 2 \\
      --data '[["Metric","Q3","Q4"], \\
               ["Revenue","$12.1M","$14.8M"], \\
               ["Margin","42%","47%"]]' \\
      --header-row

  # Plain comparison grid, no header styling
  python add_table_slide.py /tmp/deck.pptx --slide 2 \\
      --left 1 --top 1 --width 11 --height 4 \\
      --data '[["Feature","Plan A","Plan B"], \\
               ["SLA","99.9%","99.99%"], \\
               ["Support","8x5","24x7"]]'

  # Read data from JSON file (handy for long tables)
  python add_table_slide.py /tmp/deck.pptx --slide 3 \\
      --left 0.5 --top 1 --width 12 --height 5 \\
      --data-file ./roster.json --header-row --font-size 14

Data shape
----------

  JSON 2D array. Outer array = rows. Inner arrays = cells in that row.
  All rows must have the same number of columns (validated).
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
except ImportError:  # pragma: no cover - guarded by skill requires.python
    print(
        "Error: python-pptx is not installed. Run: pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(2)


_HEX_RE = re.compile(r"^[0-9A-Fa-f]{6}$")


def _hex_to_rgb(hex_str: str) -> RGBColor:
    v = hex_str.lstrip("#")
    if not _HEX_RE.match(v):
        raise ValueError(f"Color must be 6-digit hex, got {hex_str!r}")
    return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def add_table_slide(
    pptx: Path,
    *,
    slide_index: int,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    data: list[list],
    header_row: bool = False,
    font_size: int | None = None,
    header_bg: str | None = None,
    header_font_color: str | None = None,
) -> None:
    if not isinstance(data, list) or not data:
        raise ValueError("--data must be a non-empty JSON 2D array")
    if not all(isinstance(row, list) for row in data):
        raise ValueError("--data must be an array of arrays")
    n_cols = len(data[0])
    if not all(len(row) == n_cols for row in data):
        raise ValueError("All rows in --data must have the same column count")
    if n_cols == 0:
        raise ValueError("--data rows must contain at least one cell")

    prs = Presentation(pptx)
    if not 0 <= slide_index < len(prs.slides):
        raise ValueError(
            f"--slide must be in range 0..{len(prs.slides) - 1}, got {slide_index}"
        )
    slide = prs.slides[slide_index]

    n_rows = len(data)
    table_shape = slide.shapes.add_table(
        n_rows,
        n_cols,
        Inches(left_in),
        Inches(top_in),
        Inches(width_in),
        Inches(height_in),
    )
    table = table_shape.table

    hdr_bg_color = _hex_to_rgb(header_bg) if header_bg else None
    hdr_text_color = _hex_to_rgb(header_font_color) if header_font_color else None

    for r_idx, row in enumerate(data):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = "" if val is None else str(val)
            # Format the run we just created via .text
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if font_size is not None:
                        run.font.size = Pt(font_size)
                    if header_row and r_idx == 0:
                        run.font.bold = True
                        if hdr_text_color is not None:
                            run.font.color.rgb = hdr_text_color
            # Cell-level fill for header row
            if header_row and r_idx == 0 and hdr_bg_color is not None:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hdr_bg_color

    prs.save(pptx)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("pptx", type=Path, help="Path to the .pptx file")
    parser.add_argument("--slide", type=int, required=True, help="Zero-based slide index")
    parser.add_argument("--left", type=float, required=True, help="Left position (inches)")
    parser.add_argument("--top", type=float, required=True, help="Top position (inches)")
    parser.add_argument("--width", type=float, required=True, help="Table width (inches)")
    parser.add_argument("--height", type=float, required=True, help="Table height (inches)")

    data_src = parser.add_mutually_exclusive_group(required=True)
    data_src.add_argument("--data", help="JSON 2D array")
    data_src.add_argument("--data-file", type=Path, help="Path to a file containing the JSON 2D array")

    parser.add_argument("--header-row", action="store_true", help="Style row 0 as bold header")
    parser.add_argument("--font-size", type=int, default=None, help="Cell font size in points")
    parser.add_argument(
        "--header-bg",
        default=None,
        help="Header row background color, 6-digit hex like 1F2937 (no '#')",
    )
    parser.add_argument(
        "--header-font-color",
        default=None,
        help="Header row text color, 6-digit hex like FFFFFF (no '#')",
    )
    args = parser.parse_args(argv)

    if not args.pptx.is_file():
        print(f"Error: pptx not found: {args.pptx}", file=sys.stderr)
        return 1

    try:
        if args.data:
            data = json.loads(args.data)
        else:
            if not args.data_file.is_file():
                print(f"Error: --data-file not found: {args.data_file}", file=sys.stderr)
                return 1
            data = json.loads(args.data_file.read_text(encoding="utf-8"))
    except json.JSONDecodeError as e:
        print(f"Error: --data is not valid JSON: {e}", file=sys.stderr)
        return 1

    try:
        add_table_slide(
            args.pptx,
            slide_index=args.slide,
            left_in=args.left,
            top_in=args.top,
            width_in=args.width,
            height_in=args.height,
            data=data,
            header_row=args.header_row,
            font_size=args.font_size,
            header_bg=args.header_bg,
            header_font_color=args.header_font_color,
        )
    except ValueError as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1

    rows = len(data)
    cols = len(data[0]) if data else 0
    print(f"Added {rows}x{cols} table on slide {args.slide} of {args.pptx}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
