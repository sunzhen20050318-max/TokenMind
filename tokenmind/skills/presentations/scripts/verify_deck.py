#!/usr/bin/env python3
"""Audit a deck's structure and optionally assert expectations.

By default this is a reporter — exit 0 with a JSON shape summary on
stdout. Pass any ``--expect-*`` flags to turn it into a gate (non-zero
exit + a ``failures`` list in the JSON output).

Examples
--------

  # Plain JSON report
  python verify_deck.py /tmp/deck.pptx

  # Human-readable
  python verify_deck.py /tmp/deck.pptx --format text

  # Minimum slide count + every slide must have a title
  python verify_deck.py /tmp/deck.pptx \\
      --expect-min-slides 5 --expect-all-titled

  # At least one chart somewhere in the deck
  python verify_deck.py /tmp/deck.pptx --expect-charts-total ">=1"

  # No empty slides (every slide has at least one shape)
  python verify_deck.py /tmp/deck.pptx --expect-no-empty-slides

Reported fields per slide
-------------------------

  - index (zero-based)
  - title (string or null)
  - shape_count
  - chart_count
  - table_count
  - image_count
  - textbox_count
  - placeholder_count
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:  # pragma: no cover - guarded by skill requires.python
    print(
        "Error: python-pptx is not installed. Run: pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(2)


_EXPR_RE = re.compile(r"^\s*(>=|<=|=|>|<)\s*(\d+)\s*$")


def _audit_slide(slide, idx: int) -> dict[str, Any]:
    title = None
    if slide.shapes.title is not None:
        text = slide.shapes.title.text_frame.text.strip()
        title = text if text else None

    chart_count = 0
    table_count = 0
    image_count = 0
    textbox_count = 0
    placeholder_count = 0
    shape_count = 0

    for shape in slide.shapes:
        shape_count += 1
        if shape.has_chart:
            chart_count += 1
        elif shape.has_table:
            table_count += 1
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_count += 1
        elif shape.is_placeholder:
            placeholder_count += 1
        elif shape.has_text_frame:
            textbox_count += 1

    return {
        "index": idx,
        "title": title,
        "shape_count": shape_count,
        "chart_count": chart_count,
        "table_count": table_count,
        "image_count": image_count,
        "textbox_count": textbox_count,
        "placeholder_count": placeholder_count,
    }


def audit_deck(pptx: Path) -> dict[str, Any]:
    prs = Presentation(pptx)
    slides = [_audit_slide(s, i) for i, s in enumerate(prs.slides)]
    return {
        "path": str(pptx),
        "slide_count": len(slides),
        "totals": {
            "charts":      sum(s["chart_count"] for s in slides),
            "tables":      sum(s["table_count"] for s in slides),
            "images":      sum(s["image_count"] for s in slides),
            "textboxes":   sum(s["textbox_count"] for s in slides),
            "placeholders": sum(s["placeholder_count"] for s in slides),
        },
        "slides": slides,
    }


def _parse_expr(expr: str, label: str) -> tuple[str, int]:
    """Parse ``>=5`` / ``=3`` / ``<10`` into (op, n)."""
    m = _EXPR_RE.match(expr)
    if not m:
        raise ValueError(f"{label} must look like '>=5' / '=3' / '<10', got {expr!r}")
    return m.group(1), int(m.group(2))


def _cmp(actual: int, op: str, expected: int) -> bool:
    return {
        ">=": actual >= expected,
        "<=": actual <= expected,
        "=":  actual == expected,
        ">":  actual >  expected,
        "<":  actual <  expected,
    }[op]


def check_expectations(
    report: dict[str, Any],
    *,
    expect_min_slides: int | None = None,
    expect_all_titled: bool = False,
    expect_charts_total: str | None = None,
    expect_images_total: str | None = None,
    expect_tables_total: str | None = None,
    expect_no_empty_slides: bool = False,
) -> list[str]:
    failures: list[str] = []

    if expect_min_slides is not None and report["slide_count"] < expect_min_slides:
        failures.append(
            f"slide_count={report['slide_count']} < expected {expect_min_slides}"
        )

    if expect_all_titled:
        for s in report["slides"]:
            if s["title"] is None:
                failures.append(f"slide {s['index']} has no title")

    for label, key, expr in [
        ("--expect-charts-total", "charts", expect_charts_total),
        ("--expect-images-total", "images", expect_images_total),
        ("--expect-tables-total", "tables", expect_tables_total),
    ]:
        if expr is None:
            continue
        op, n = _parse_expr(expr, label)
        actual = report["totals"][key]
        if not _cmp(actual, op, n):
            failures.append(f"total {key}={actual} fails {op}{n}")

    if expect_no_empty_slides:
        for s in report["slides"]:
            if s["shape_count"] == 0:
                failures.append(f"slide {s['index']} is empty")

    return failures


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("pptx", type=Path, help="Path to the .pptx file")
    parser.add_argument("--format", choices=["json", "text"], default="json", help="Output format (default: json)")
    parser.add_argument("--expect-min-slides", type=int, default=None, help="Fail if slide_count < N")
    parser.add_argument("--expect-all-titled", action="store_true", help="Every slide must have a non-empty title")
    parser.add_argument("--expect-charts-total", default=None, help="Comparison expression, e.g. '>=1'")
    parser.add_argument("--expect-images-total", default=None, help="Comparison expression, e.g. '>=2'")
    parser.add_argument("--expect-tables-total", default=None, help="Comparison expression, e.g. '<5'")
    parser.add_argument("--expect-no-empty-slides", action="store_true", help="Every slide must contain ≥ 1 shape")
    args = parser.parse_args(argv)

    if not args.pptx.is_file():
        print(f"Error: pptx not found: {args.pptx}", file=sys.stderr)
        return 1

    try:
        report = audit_deck(args.pptx)
    except Exception as e:  # noqa: BLE001
        print(f"Error: failed to load deck: {e}", file=sys.stderr)
        return 1

    try:
        failures = check_expectations(
            report,
            expect_min_slides=args.expect_min_slides,
            expect_all_titled=args.expect_all_titled,
            expect_charts_total=args.expect_charts_total,
            expect_images_total=args.expect_images_total,
            expect_tables_total=args.expect_tables_total,
            expect_no_empty_slides=args.expect_no_empty_slides,
        )
    except ValueError as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1

    if args.format == "json":
        print(json.dumps({"report": report, "failures": failures}, indent=2, ensure_ascii=False))
    else:
        print(f"Deck: {report['path']}")
        print(f"  slides: {report['slide_count']}")
        print(f"  totals: {report['totals']}")
        for s in report["slides"]:
            print(
                f"  - {s['index']}: title={s['title']!r} "
                f"charts={s['chart_count']} tables={s['table_count']} "
                f"images={s['image_count']} textboxes={s['textbox_count']} "
                f"placeholders={s['placeholder_count']}"
            )
        if failures:
            print("Failures:")
            for f in failures:
                print(f"  - {f}")
        else:
            print("OK")

    return 1 if failures else 0


if __name__ == "__main__":
    sys.exit(main())
