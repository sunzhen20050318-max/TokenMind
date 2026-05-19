#!/usr/bin/env python3
"""Add a slide to an existing .pptx using one of the default layouts.

Default python-pptx layout map (zero-based, names may vary by template):

  0  Title Slide              (title + subtitle)
  1  Title and Content        (title + body)
  2  Section Header
  3  Two Content              (title + two body columns)
  4  Comparison
  5  Title Only               (title, no body — good for figure slides)
  6  Blank                    (empty canvas — best for free-form composition)
  7  Content with Caption
  8  Picture with Caption

Use ``--list-layouts`` to dump the actual layout names present in the deck.

Examples
--------

  # Append a "Title Only" slide ready for free-form composition
  python add_slide.py /tmp/deck.pptx --layout 5

  # Append a Title and Content slide, fill the title and bullet body
  python add_slide.py /tmp/deck.pptx --layout 1 \\
      --title "Key Findings" \\
      --bullets '["Revenue up 12% YoY","Margin expanded 180 bps","Churn flat"]'

  # Insert a blank slide as slide #2 (zero-based position)
  python add_slide.py /tmp/deck.pptx --layout 6 --position 1

  # Inspect what layouts the deck actually has
  python add_slide.py /tmp/deck.pptx --list-layouts
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover - guarded by skill requires.python
    print(
        "Error: python-pptx is not installed. Run: pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(2)


def list_layouts(workbook: Path) -> list[dict]:
    prs = Presentation(workbook)
    return [
        {"index": i, "name": layout.name}
        for i, layout in enumerate(prs.slide_layouts)
    ]


def add_slide(
    pptx: Path,
    *,
    layout_index: int,
    title: str | None = None,
    subtitle: str | None = None,
    bullets: list[str] | None = None,
    position: int | None = None,
) -> int:
    """Append (or insert) a slide. Returns the final slide count."""
    prs = Presentation(pptx)

    if not 0 <= layout_index < len(prs.slide_layouts):
        raise ValueError(
            f"--layout must be in range 0..{len(prs.slide_layouts) - 1}, got {layout_index}"
        )

    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)

    if title is not None:
        if slide.shapes.title is None:
            raise ValueError(f"Layout {layout_index} has no title placeholder")
        slide.shapes.title.text = title

    if subtitle is not None:
        # Subtitle is usually placeholder idx=1 on Title Slide layouts
        subtitle_ph = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1 and ph != slide.shapes.title:
                subtitle_ph = ph
                break
        if subtitle_ph is None:
            raise ValueError(f"Layout {layout_index} has no subtitle/body placeholder")
        subtitle_ph.text = subtitle

    if bullets:
        # Find a body placeholder (any non-title placeholder)
        body_ph = None
        for ph in slide.placeholders:
            if ph == slide.shapes.title:
                continue
            body_ph = ph
            break
        if body_ph is None:
            raise ValueError(
                f"Layout {layout_index} has no body placeholder to receive bullets"
            )
        tf = body_ph.text_frame
        tf.text = bullets[0]
        for line in bullets[1:]:
            p = tf.add_paragraph()
            p.text = line

    # Move slide if a position was requested (python-pptx adds to end).
    if position is not None:
        if not 0 <= position < len(prs.slides):
            raise ValueError(
                f"--position must be in range 0..{len(prs.slides) - 1}, got {position}"
            )
        # XML reorder — python-pptx doesn't expose a high-level move, but the
        # sldIdLst child order is the canonical source of truth.
        sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
        slides = list(sldIdLst)
        new_slide_xml = slides[-1]  # the one we just appended
        sldIdLst.remove(new_slide_xml)
        sldIdLst.insert(position, new_slide_xml)

    prs.save(pptx)
    return len(prs.slides)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("pptx", type=Path, help="Path to the .pptx file")
    parser.add_argument("--layout", type=int, default=None, help="Layout index (see --list-layouts)")
    parser.add_argument("--title", default=None, help="Title text")
    parser.add_argument("--subtitle", default=None, help="Subtitle text (Title Slide layout)")
    parser.add_argument(
        "--bullets",
        default=None,
        help='JSON array of bullet strings: \'["one","two","three"]\'',
    )
    parser.add_argument(
        "--position",
        type=int,
        default=None,
        help="Zero-based insertion index (default: append)",
    )
    parser.add_argument(
        "--list-layouts",
        action="store_true",
        help="Print available layouts in this deck and exit",
    )
    args = parser.parse_args(argv)

    if not args.pptx.is_file():
        print(f"Error: pptx not found: {args.pptx}", file=sys.stderr)
        return 1

    if args.list_layouts:
        for entry in list_layouts(args.pptx):
            print(f"{entry['index']:>2}  {entry['name']}")
        return 0

    if args.layout is None:
        print("Error: --layout is required (use --list-layouts to inspect)", file=sys.stderr)
        return 1

    bullets = None
    if args.bullets:
        try:
            bullets = json.loads(args.bullets)
        except json.JSONDecodeError as e:
            print(f"Error: --bullets must be valid JSON: {e}", file=sys.stderr)
            return 1
        if not isinstance(bullets, list) or not all(isinstance(b, str) for b in bullets):
            print("Error: --bullets must be a JSON array of strings", file=sys.stderr)
            return 1

    try:
        total = add_slide(
            args.pptx,
            layout_index=args.layout,
            title=args.title,
            subtitle=args.subtitle,
            bullets=bullets,
            position=args.position,
        )
    except ValueError as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1

    print(f"Added slide (layout {args.layout}) to {args.pptx} (total slides: {total})")
    return 0


if __name__ == "__main__":
    sys.exit(main())
