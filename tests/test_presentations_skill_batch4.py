"""Tests for presentations skill batch 4: verify_deck / apply_industry_profile."""
from __future__ import annotations

import importlib.util
import json
import subprocess
import sys
from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402

SKILL_DIR = Path(__file__).resolve().parent.parent / "tokenmind" / "skills" / "presentations" / "scripts"


def _load_script(path: Path):
    spec = importlib.util.spec_from_file_location(path.stem, path)
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    sys.modules[path.stem] = module
    spec.loader.exec_module(module)
    return module


@pytest.fixture(scope="module")
def verify_deck():
    return _load_script(SKILL_DIR / "verify_deck.py")


@pytest.fixture(scope="module")
def apply_industry_profile():
    return _load_script(SKILL_DIR / "apply_industry_profile.py")


def _make_rich_deck(tmp_path: Path) -> Path:
    """Deck with: titled slide, untitled slide, slide w/ chart + table + image."""
    out = tmp_path / "rich.pptx"
    prs = Presentation()
    from pptx.util import Inches

    # Slide 0: Title + Content (titled)
    s0 = prs.slides.add_slide(prs.slide_layouts[1])
    s0.shapes.title.text = "Agenda"
    body = next(ph for ph in s0.placeholders if ph != s0.shapes.title)
    body.text_frame.text = "Intro"

    # Slide 1: Title Only but with empty title (counts as untitled in audit)
    s1 = prs.slides.add_slide(prs.slide_layouts[5])

    # Slide 2: Blank with table + textbox + chart
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
    tb.text_frame.text = "Body"
    tbl = s2.shapes.add_table(2, 2, Inches(0.5), Inches(2), Inches(5), Inches(2))
    tbl.table.cell(0, 0).text = "A"
    tbl.table.cell(0, 1).text = "B"

    # Chart on s2
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    cd = CategoryChartData()
    cd.categories = ["Q1", "Q2"]
    cd.add_series("Revenue", [100, 200])
    s2.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(6), Inches(2), Inches(6), Inches(4),
        cd,
    )

    prs.save(out)
    return out


# --- verify_deck ------------------------------------------------------------


def test_verify_reports_per_slide_counts(tmp_path, verify_deck):
    deck = _make_rich_deck(tmp_path)
    report = verify_deck.audit_deck(deck)
    assert report["slide_count"] == 3
    assert report["totals"]["charts"] == 1
    assert report["totals"]["tables"] == 1
    assert report["totals"]["textboxes"] == 1

    s0 = report["slides"][0]
    assert s0["title"] == "Agenda"
    assert s0["placeholder_count"] >= 1  # title + body placeholders

    s1 = report["slides"][1]
    assert s1["title"] is None  # empty title text

    s2 = report["slides"][2]
    assert s2["chart_count"] == 1
    assert s2["table_count"] == 1
    assert s2["textbox_count"] == 1


def test_verify_expect_min_slides(tmp_path, verify_deck):
    deck = _make_rich_deck(tmp_path)
    report = verify_deck.audit_deck(deck)
    assert verify_deck.check_expectations(report, expect_min_slides=3) == []
    failures = verify_deck.check_expectations(report, expect_min_slides=10)
    assert any("slide_count" in f for f in failures)


def test_verify_expect_all_titled_catches_untitled(tmp_path, verify_deck):
    deck = _make_rich_deck(tmp_path)
    report = verify_deck.audit_deck(deck)
    failures = verify_deck.check_expectations(report, expect_all_titled=True)
    # slide 1 has no title, slide 2 too (blank layout)
    assert len(failures) >= 1
    assert any("slide 1" in f for f in failures)


def test_verify_expect_charts_expression(tmp_path, verify_deck):
    deck = _make_rich_deck(tmp_path)
    report = verify_deck.audit_deck(deck)
    assert verify_deck.check_expectations(report, expect_charts_total=">=1") == []
    failures = verify_deck.check_expectations(report, expect_charts_total=">=5")
    assert any("charts" in f for f in failures)


def test_verify_expect_no_empty_slides(tmp_path, verify_deck):
    deck = _make_rich_deck(tmp_path)
    report = verify_deck.audit_deck(deck)
    # Slide 1 is Title Only with empty title — but has a title placeholder shape
    failures = verify_deck.check_expectations(report, expect_no_empty_slides=True)
    # Title-Only layout still has 1 shape (the placeholder), so this passes.
    # Test with a truly empty Blank slide:
    prs = Presentation(deck)
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    _ = blank
    prs.save(deck)
    report2 = verify_deck.audit_deck(deck)
    failures2 = verify_deck.check_expectations(report2, expect_no_empty_slides=True)
    assert any("empty" in f for f in failures2)


def test_verify_bad_expression(tmp_path, verify_deck):
    deck = _make_rich_deck(tmp_path)
    report = verify_deck.audit_deck(deck)
    with pytest.raises(ValueError, match="must look like"):
        verify_deck.check_expectations(report, expect_charts_total="bogus")


def test_verify_cli_json(tmp_path):
    deck = _make_rich_deck(tmp_path)
    r = subprocess.run(
        [sys.executable, str(SKILL_DIR / "verify_deck.py"), str(deck)],
        capture_output=True, text=True,
    )
    assert r.returncode == 0, r.stderr
    payload = json.loads(r.stdout)
    assert payload["report"]["slide_count"] == 3
    assert payload["failures"] == []


def test_verify_cli_nonzero_on_failure(tmp_path):
    deck = _make_rich_deck(tmp_path)
    r = subprocess.run(
        [
            sys.executable, str(SKILL_DIR / "verify_deck.py"), str(deck),
            "--expect-min-slides", "10",
        ],
        capture_output=True, text=True,
    )
    assert r.returncode == 1
    payload = json.loads(r.stdout)
    assert len(payload["failures"]) >= 1


# --- apply_industry_profile -------------------------------------------------


def test_seven_profiles_available(apply_industry_profile):
    profiles = list(apply_industry_profile.PROFILES.keys())
    assert len(profiles) == 7
    assert set(profiles) == {
        "consumer-retail",
        "engineering-platform",
        "finance-ir",
        "gtm-growth",
        "product-platform",
        "appendix-heavy",
        "template-and-edit",
    }


def test_get_profile_returns_copy(apply_industry_profile):
    a = apply_industry_profile.get_profile("finance-ir")
    a["accent_color"] = "FF0000"
    b = apply_industry_profile.get_profile("finance-ir")
    assert b["accent_color"] == "C19A28"  # not mutated


def test_get_profile_unknown_raises(apply_industry_profile):
    with pytest.raises(ValueError, match="Unknown profile"):
        apply_industry_profile.get_profile("not-a-profile")


def test_apply_profile_changes_colors(tmp_path, apply_industry_profile):
    deck = _make_rich_deck(tmp_path)
    spec = apply_industry_profile.apply_profile(deck, "finance-ir")
    assert spec["heading_color"] == "0B2545"

    prs = Presentation(deck)
    title_run = prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0]
    assert str(title_run.font.color.rgb) == "0B2545"


def test_apply_profile_with_override(tmp_path, apply_industry_profile):
    deck = _make_rich_deck(tmp_path)
    spec = apply_industry_profile.apply_profile(
        deck, "product-platform", overrides={"accent_color": "8B5CF6"},
    )
    assert spec["accent_color"] == "8B5CF6"
    # Other keys unchanged from profile default
    assert spec["heading_color"] == "000000"


def test_cli_list(apply_industry_profile):
    r = subprocess.run(
        [sys.executable, str(SKILL_DIR / "apply_industry_profile.py"), "--list"],
        capture_output=True, text=True,
    )
    assert r.returncode == 0
    listed = [line.strip() for line in r.stdout.splitlines() if line.strip()]
    assert "finance-ir" in listed
    assert len(listed) == 7


def test_cli_dump(tmp_path):
    r = subprocess.run(
        [
            sys.executable, str(SKILL_DIR / "apply_industry_profile.py"),
            "--profile", "engineering-platform", "--dump",
        ],
        capture_output=True, text=True,
    )
    assert r.returncode == 0
    spec = json.loads(r.stdout)
    assert spec["accent_color"] == "06B6D4"


def test_cli_dump_with_override(tmp_path):
    r = subprocess.run(
        [
            sys.executable, str(SKILL_DIR / "apply_industry_profile.py"),
            "--profile", "engineering-platform",
            "--dump",
            "--override", "accent_color=FF00FF",
        ],
        capture_output=True, text=True,
    )
    assert r.returncode == 0
    spec = json.loads(r.stdout)
    assert spec["accent_color"] == "FF00FF"


def test_cli_apply_smoke(tmp_path):
    deck = _make_rich_deck(tmp_path)
    r = subprocess.run(
        [
            sys.executable, str(SKILL_DIR / "apply_industry_profile.py"),
            str(deck),
            "--profile", "consumer-retail",
        ],
        capture_output=True, text=True,
    )
    assert r.returncode == 0, r.stderr

    prs = Presentation(deck)
    title_run = prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0]
    assert str(title_run.font.color.rgb) == "1F1F1F"


def test_cli_bad_override_rejected(tmp_path):
    r = subprocess.run(
        [
            sys.executable, str(SKILL_DIR / "apply_industry_profile.py"),
            "--profile", "finance-ir", "--dump",
            "--override", "no-equals-sign",
        ],
        capture_output=True, text=True,
    )
    assert r.returncode == 1
    assert "key=value" in r.stderr
