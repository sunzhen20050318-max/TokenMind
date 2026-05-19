"""Tests for presentations skill batch 1: build_pptx / add_slide / add_text_box."""
from __future__ import annotations

import importlib.util
import json
import subprocess
import sys
from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

SKILL_DIR = Path(__file__).resolve().parent.parent / "tokenmind" / "skills" / "presentations" / "scripts"


def _load_script(path: Path):
    spec = importlib.util.spec_from_file_location(path.stem, path)
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    sys.modules[path.stem] = module
    spec.loader.exec_module(module)
    return module


@pytest.fixture(scope="module")
def build_pptx():
    return _load_script(SKILL_DIR / "build_pptx.py")


@pytest.fixture(scope="module")
def add_slide():
    return _load_script(SKILL_DIR / "add_slide.py")


@pytest.fixture(scope="module")
def add_text_box():
    return _load_script(SKILL_DIR / "add_text_box.py")


# --- build_pptx --------------------------------------------------------------


def test_build_pptx_default_16x9(tmp_path, build_pptx):
    out = tmp_path / "deck.pptx"
    result = build_pptx.build_pptx(out)
    assert result == out
    assert out.is_file()
    prs = Presentation(out)
    assert prs.slide_width == Inches(13.333)
    assert prs.slide_height == Inches(7.5)
    assert len(prs.slides) == 0


def test_build_pptx_4x3(tmp_path, build_pptx):
    out = tmp_path / "deck.pptx"
    build_pptx.build_pptx(out, size="4:3")
    prs = Presentation(out)
    assert prs.slide_width == Inches(10.0)
    assert prs.slide_height == Inches(7.5)


def test_build_pptx_custom_size(tmp_path, build_pptx):
    out = tmp_path / "deck.pptx"
    build_pptx.build_pptx(out, width_in=11.69, height_in=8.27)
    prs = Presentation(out)
    assert prs.slide_width == Inches(11.69)
    assert prs.slide_height == Inches(8.27)


def test_build_pptx_partial_custom_rejected(tmp_path, build_pptx):
    out = tmp_path / "deck.pptx"
    with pytest.raises(ValueError, match="both"):
        build_pptx.build_pptx(out, width_in=11.69)  # missing height


def test_build_pptx_cover_title(tmp_path, build_pptx):
    out = tmp_path / "deck.pptx"
    build_pptx.build_pptx(out, cover_title="Q3 Results")
    prs = Presentation(out)
    assert len(prs.slides) == 1
    assert prs.slides[0].shapes.title.text == "Q3 Results"


def test_build_pptx_refuses_clobber(tmp_path, build_pptx):
    out = tmp_path / "deck.pptx"
    build_pptx.build_pptx(out)
    with pytest.raises(FileExistsError):
        build_pptx.build_pptx(out)


def test_build_pptx_overwrite(tmp_path, build_pptx):
    out = tmp_path / "deck.pptx"
    build_pptx.build_pptx(out)
    build_pptx.build_pptx(out, cover_title="Replaced", overwrite=True)
    prs = Presentation(out)
    assert prs.slides[0].shapes.title.text == "Replaced"


def test_build_pptx_creates_parent_dirs(tmp_path, build_pptx):
    out = tmp_path / "nested" / "deep" / "deck.pptx"
    build_pptx.build_pptx(out)
    assert out.is_file()


# --- add_slide ---------------------------------------------------------------


def _make_deck(tmp_path: Path) -> Path:
    out = tmp_path / "deck.pptx"
    Presentation().save(out)
    return out


def test_add_slide_title_and_content(tmp_path, add_slide):
    deck = _make_deck(tmp_path)
    total = add_slide.add_slide(
        deck,
        layout_index=1,
        title="Key Findings",
        bullets=["Revenue up 12% YoY", "Margin expanded 180 bps", "Churn flat"],
    )
    assert total == 1
    prs = Presentation(deck)
    slide = prs.slides[0]
    assert slide.shapes.title.text == "Key Findings"
    # Body placeholder should hold bullets
    body = None
    for ph in slide.placeholders:
        if ph != slide.shapes.title:
            body = ph
            break
    assert body is not None
    text = body.text_frame.text.split("\n")
    assert "Revenue up 12% YoY" in text
    assert "Margin expanded 180 bps" in text


def test_add_slide_title_only_layout(tmp_path, add_slide):
    deck = _make_deck(tmp_path)
    add_slide.add_slide(deck, layout_index=5, title="Figure 1")
    prs = Presentation(deck)
    assert prs.slides[0].shapes.title.text == "Figure 1"


def test_add_slide_blank_layout(tmp_path, add_slide):
    deck = _make_deck(tmp_path)
    total = add_slide.add_slide(deck, layout_index=6)
    assert total == 1
    prs = Presentation(deck)
    # Blank layout has no title placeholder
    assert prs.slides[0].shapes.title is None


def test_add_slide_position_insertion(tmp_path, add_slide):
    deck = _make_deck(tmp_path)
    add_slide.add_slide(deck, layout_index=1, title="A")
    add_slide.add_slide(deck, layout_index=1, title="B")
    # Insert C at position 1 (between A and B)
    add_slide.add_slide(deck, layout_index=1, title="C", position=1)
    prs = Presentation(deck)
    titles = [s.shapes.title.text for s in prs.slides]
    assert titles == ["A", "C", "B"]


def test_add_slide_invalid_layout(tmp_path, add_slide):
    deck = _make_deck(tmp_path)
    with pytest.raises(ValueError, match="--layout"):
        add_slide.add_slide(deck, layout_index=99)


def test_add_slide_list_layouts(tmp_path, add_slide):
    deck = _make_deck(tmp_path)
    layouts = add_slide.list_layouts(deck)
    assert len(layouts) >= 9
    names = {entry["name"] for entry in layouts}
    # Some layout naming differs by template; just sanity-check the count
    assert any("Title" in n for n in names)


# --- add_text_box ------------------------------------------------------------


def test_add_text_box_basic(tmp_path, add_slide, add_text_box):
    deck = _make_deck(tmp_path)
    add_slide.add_slide(deck, layout_index=6)  # Blank
    add_text_box.add_text_box(
        deck,
        slide_index=0,
        left_in=0.5, top_in=0.4, width_in=9, height_in=1,
        text="Quarterly Results",
        font_size=36, bold=True,
    )
    prs = Presentation(deck)
    slide = prs.slides[0]
    text_shapes = [s for s in slide.shapes if s.has_text_frame]
    assert len(text_shapes) == 1
    tf = text_shapes[0].text_frame
    assert tf.paragraphs[0].text == "Quarterly Results"
    run = tf.paragraphs[0].runs[0]
    assert run.font.bold is True
    # Pt value comparison
    from pptx.util import Pt
    assert run.font.size == Pt(36)


def test_add_text_box_multiline_split(tmp_path, add_slide, add_text_box):
    deck = _make_deck(tmp_path)
    add_slide.add_slide(deck, layout_index=6)
    add_text_box.add_text_box(
        deck,
        slide_index=0,
        left_in=0.5, top_in=1, width_in=9, height_in=4,
        text="Revenue: $12.4M\nGrowth: +18%\nChurn: 2.1%",
        font_size=22,
    )
    prs = Presentation(deck)
    tf = next(s for s in prs.slides[0].shapes if s.has_text_frame).text_frame
    lines = [p.text for p in tf.paragraphs]
    assert lines == ["Revenue: $12.4M", "Growth: +18%", "Churn: 2.1%"]


def test_add_text_box_font_color(tmp_path, add_slide, add_text_box):
    deck = _make_deck(tmp_path)
    add_slide.add_slide(deck, layout_index=6)
    add_text_box.add_text_box(
        deck,
        slide_index=0,
        left_in=0.5, top_in=1, width_in=9, height_in=1,
        text="Gray", font_size=18, font_color="4B5563",
    )
    prs = Presentation(deck)
    tf = next(s for s in prs.slides[0].shapes if s.has_text_frame).text_frame
    color = tf.paragraphs[0].runs[0].font.color.rgb
    assert str(color) == "4B5563"


def test_add_text_box_bad_hex_rejected(tmp_path, add_slide, add_text_box):
    deck = _make_deck(tmp_path)
    add_slide.add_slide(deck, layout_index=6)
    with pytest.raises(ValueError, match="hex"):
        add_text_box.add_text_box(
            deck, slide_index=0, left_in=0, top_in=0, width_in=1, height_in=1,
            text="x", font_color="ZZZZZZ",
        )


def test_add_text_box_bad_align_rejected(tmp_path, add_slide, add_text_box):
    deck = _make_deck(tmp_path)
    add_slide.add_slide(deck, layout_index=6)
    with pytest.raises(ValueError, match="--align"):
        add_text_box.add_text_box(
            deck, slide_index=0, left_in=0, top_in=0, width_in=1, height_in=1,
            text="x", align="middle",
        )


def test_add_text_box_slide_out_of_range(tmp_path, add_text_box):
    deck = _make_deck(tmp_path)
    # No slides yet → slide 0 is invalid
    with pytest.raises(ValueError, match="--slide"):
        add_text_box.add_text_box(
            deck, slide_index=0, left_in=0, top_in=0, width_in=1, height_in=1,
            text="x",
        )


# --- end-to-end CLI smoke ----------------------------------------------------


def test_cli_smoke(tmp_path):
    deck = tmp_path / "deck.pptx"

    r = subprocess.run(
        [sys.executable, str(SKILL_DIR / "build_pptx.py"), "--out", str(deck), "--cover-title", "Hello"],
        capture_output=True, text=True,
    )
    assert r.returncode == 0, r.stderr

    r = subprocess.run(
        [
            sys.executable, str(SKILL_DIR / "add_slide.py"), str(deck),
            "--layout", "1",
            "--title", "Agenda",
            "--bullets", json.dumps(["Intro", "Findings", "Next steps"]),
        ],
        capture_output=True, text=True,
    )
    assert r.returncode == 0, r.stderr

    r = subprocess.run(
        [
            sys.executable, str(SKILL_DIR / "add_slide.py"), str(deck),
            "--layout", "6",  # Blank
        ],
        capture_output=True, text=True,
    )
    assert r.returncode == 0, r.stderr

    r = subprocess.run(
        [
            sys.executable, str(SKILL_DIR / "add_text_box.py"), str(deck),
            "--slide", "2",
            "--left", "1.0", "--top", "1.0", "--width", "10", "--height", "4",
            "--text", "Body text on a blank slide",
            "--font-size", "24",
        ],
        capture_output=True, text=True,
    )
    assert r.returncode == 0, r.stderr

    prs = Presentation(deck)
    assert len(prs.slides) == 3
    assert prs.slides[0].shapes.title.text == "Hello"
    assert prs.slides[1].shapes.title.text == "Agenda"
