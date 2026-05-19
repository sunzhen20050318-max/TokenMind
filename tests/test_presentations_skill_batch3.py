"""Tests for presentations skill batch 3: apply_design_system / render_deck."""
from __future__ import annotations

import importlib.util
import json
import shutil
import subprocess
import sys
from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402

SKILL_DIR = Path(__file__).resolve().parent.parent / "tokenmind" / "skills" / "presentations" / "scripts"


def _load_script(path: Path):
    spec = importlib.util.spec_from_file_location(path.stem, path)
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    sys.modules[path.stem] = module
    spec.loader.exec_module(module)
    return module


@pytest.fixture(scope="module")
def apply_design_system():
    return _load_script(SKILL_DIR / "apply_design_system.py")


@pytest.fixture(scope="module")
def render_deck():
    return _load_script(SKILL_DIR / "render_deck.py")


def _make_styled_deck(tmp_path: Path) -> Path:
    """Deck with a Title+Content slide, a Blank slide w/ textbox, and a table."""
    out = tmp_path / "deck.pptx"
    prs = Presentation()

    # Slide 0: Title and Content
    s0 = prs.slides.add_slide(prs.slide_layouts[1])
    s0.shapes.title.text = "Quarterly Review"
    body = next(ph for ph in s0.placeholders if ph != s0.shapes.title)
    body.text_frame.text = "Revenue is up"

    # Slide 1: Blank + textbox + table
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    from pptx.util import Inches
    tb = s1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    tb.text_frame.text = "Body text"
    tbl = s1.shapes.add_table(2, 2, Inches(0.5), Inches(2), Inches(8), Inches(2))
    tbl.table.cell(0, 0).text = "Header A"
    tbl.table.cell(0, 1).text = "Header B"
    tbl.table.cell(1, 0).text = "Value 1"
    tbl.table.cell(1, 1).text = "Value 2"

    prs.save(out)
    return out


# --- apply_design_system ----------------------------------------------------


def test_design_applies_heading_and_body_colors(tmp_path, apply_design_system):
    deck = _make_styled_deck(tmp_path)
    counts = apply_design_system.apply_design_system(
        deck,
        {
            "heading_color": "111827",
            "body_color": "4B5563",
            "heading_font": "Inter",
            "body_font": "Inter",
        },
    )
    assert counts["slides"] == 2
    assert counts["titles"] >= 1
    assert counts["bodies"] >= 2  # body placeholder + textbox

    prs = Presentation(deck)
    title_run = prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0]
    assert str(title_run.font.color.rgb) == "111827"
    assert title_run.font.name == "Inter"

    # Body placeholder on slide 0
    body_ph = next(ph for ph in prs.slides[0].placeholders if ph != prs.slides[0].shapes.title)
    body_run = body_ph.text_frame.paragraphs[0].runs[0]
    assert str(body_run.font.color.rgb) == "4B5563"


def test_design_applies_to_table_header_vs_body(tmp_path, apply_design_system):
    deck = _make_styled_deck(tmp_path)
    apply_design_system.apply_design_system(
        deck,
        {"heading_color": "FF0000", "body_color": "00FF00"},
    )
    prs = Presentation(deck)
    tbl = next(s for s in prs.slides[1].shapes if s.has_table).table
    hdr_run = tbl.cell(0, 0).text_frame.paragraphs[0].runs[0]
    body_run = tbl.cell(1, 0).text_frame.paragraphs[0].runs[0]
    assert str(hdr_run.font.color.rgb) == "FF0000"
    assert str(body_run.font.color.rgb) == "00FF00"


def test_design_sets_slide_background(tmp_path, apply_design_system):
    deck = _make_styled_deck(tmp_path)
    counts = apply_design_system.apply_design_system(deck, {"slide_bg": "F8FAFC"})
    assert counts["backgrounds"] == 2  # every slide painted
    prs = Presentation(deck)
    bg = prs.slides[0].background.fill.fore_color.rgb
    assert str(bg) == "F8FAFC"


def test_design_unknown_key_rejected(tmp_path, apply_design_system):
    deck = _make_styled_deck(tmp_path)
    with pytest.raises(ValueError, match="Unknown design-system"):
        apply_design_system.apply_design_system(deck, {"font": "Arial"})  # typo


def test_design_bad_hex_rejected(tmp_path, apply_design_system):
    deck = _make_styled_deck(tmp_path)
    with pytest.raises(ValueError, match="hex"):
        apply_design_system.apply_design_system(deck, {"heading_color": "ZZZZZZ"})


def test_design_partial_spec_leaves_other_aspects_alone(tmp_path, apply_design_system):
    deck = _make_styled_deck(tmp_path)
    prs_before = Presentation(deck)
    # Capture original textbox content
    original_text = prs_before.slides[1].shapes[0].text_frame.text

    # Only change body_color — fonts and titles untouched
    apply_design_system.apply_design_system(deck, {"body_color": "AAAAAA"})

    prs_after = Presentation(deck)
    # Same textbox content
    assert prs_after.slides[1].shapes[0].text_frame.text == original_text
    # Slide count unchanged
    assert len(prs_after.slides) == 2


def test_design_empty_spec_is_noop(tmp_path, apply_design_system):
    deck = _make_styled_deck(tmp_path)
    counts = apply_design_system.apply_design_system(deck, {})
    assert counts["slides"] == 2
    assert counts["backgrounds"] == 0


def test_design_cli_smoke(tmp_path):
    deck = _make_styled_deck(tmp_path)
    spec_path = tmp_path / "brand.json"
    spec_path.write_text(
        json.dumps({"heading_color": "111827", "body_color": "4B5563"}),
        encoding="utf-8",
    )
    r = subprocess.run(
        [
            sys.executable, str(SKILL_DIR / "apply_design_system.py"), str(deck),
            "--design-file", str(spec_path),
        ],
        capture_output=True, text=True,
    )
    assert r.returncode == 0, r.stderr


# --- render_deck ------------------------------------------------------------


def _soffice_available() -> bool:
    if shutil.which("soffice"):
        return True
    if Path("/Applications/LibreOffice.app/Contents/MacOS/soffice").is_file():
        return True
    return False


def _pdf2image_works() -> bool:
    """Try a no-op import + check poppler binary."""
    try:
        from pdf2image import pdfinfo_from_path  # noqa: F401
    except ImportError:
        return False
    return shutil.which("pdftoppm") is not None or shutil.which("pdfinfo") is not None


@pytest.mark.skipif(not _soffice_available(), reason="soffice (LibreOffice) not installed")
@pytest.mark.skipif(not _pdf2image_works(), reason="poppler / pdf2image not available")
def test_render_deck_produces_png_per_slide(tmp_path, render_deck):
    deck = _make_styled_deck(tmp_path)
    out_dir = tmp_path / "render"
    result = render_deck.render_deck(deck, out_dir, dpi=72)  # low dpi for speed
    assert result["slide_count"] == 2
    assert len(result["png_paths"]) == 2
    for p in result["png_paths"]:
        assert Path(p).is_file()
        assert Path(p).stat().st_size > 0


@pytest.mark.skipif(not _soffice_available(), reason="soffice (LibreOffice) not installed")
@pytest.mark.skipif(not _pdf2image_works(), reason="poppler / pdf2image not available")
def test_render_deck_emit_pdf(tmp_path, render_deck):
    deck = _make_styled_deck(tmp_path)
    out_dir = tmp_path / "render"
    result = render_deck.render_deck(deck, out_dir, dpi=72, emit_pdf=True)
    assert result["pdf"] is not None
    assert Path(result["pdf"]).is_file()


def test_render_deck_missing_pptx_raises(tmp_path, render_deck):
    with pytest.raises(FileNotFoundError):
        render_deck.render_deck(tmp_path / "nope.pptx", tmp_path / "out")


def test_render_deck_no_soffice_helpful_error(tmp_path, monkeypatch, render_deck):
    """When soffice is genuinely missing, the error message should name it."""
    deck = _make_styled_deck(tmp_path)

    # Substitute the soffice resolver with one that always fails. This is the
    # explicit seam in render_deck; mocking shutil.which / Path.is_file
    # globally would also have to fight pytest's own filesystem use.
    def _missing(*_args, **_kwargs):
        raise FileNotFoundError(
            "soffice (LibreOffice) not found. Install via 'brew install libreoffice' or 'apt install libreoffice'."
        )

    monkeypatch.setattr(render_deck, "_ensure_soffice", _missing)

    with pytest.raises((FileNotFoundError, RuntimeError)) as excinfo:
        render_deck.render_deck(deck, tmp_path / "out")
    assert "soffice" in str(excinfo.value).lower() or "libreoffice" in str(excinfo.value).lower()
