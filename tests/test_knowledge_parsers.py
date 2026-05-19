"""Unit tests for the rich knowledge-base document parsers.

We build the test fixtures (DOCX / PPTX / XLSX) in-memory with the same
libraries the parser uses so we don't need to ship binary artifacts in
the repo. The VLM path is exercised by monkey-patching the captioning
helper instead of calling a real API.
"""

from __future__ import annotations

from pathlib import Path
from unittest.mock import patch

import pytest

from tokenmind.knowledge import parsers


# ---------------------------------------------------------------------------
# Fixture builders
# ---------------------------------------------------------------------------

def _build_docx(path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("Quarterly Report", level=1)
    doc.add_paragraph("Revenue grew steadily over the past quarter.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "ARR"
    table.cell(1, 1).text = "$12M"
    doc.add_paragraph(
        "After the table we keep talking about projections and pipeline "
        "growth so this paragraph clearly carries more than the bare "
        "table — useful for verifying ordering."
    )
    doc.save(str(path))


def _build_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    layout = prs.slide_layouts[5]  # Title Only
    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Intro"
    tb = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
    tb.text_frame.text = "Welcome to the deck."

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Numbers"
    tb2 = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
    tb2.text_frame.text = "Revenue doubled YoY."
    prs.save(str(path))


def _build_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Q1"
    ws.append(["Metric", "Value"])
    ws.append(["ARR", "$12M"])
    ws.append(["Customers", 132])

    ws2 = wb.create_sheet("Q2")
    ws2.append(["Metric", "Value"])
    ws2.append(["ARR", "$15M"])
    wb.save(str(path))


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def test_parse_docx_preserves_heading_and_table(tmp_path):
    path = tmp_path / "report.docx"
    _build_docx(path)

    doc = parsers.parse_docx(path)
    text = doc.as_text()

    assert doc.file_type == "docx"
    assert len(doc.pages) >= 1
    assert "# Quarterly Report" in text
    # Table rendered as "cell | cell" rows
    assert "Metric | Value" in text
    assert "ARR | $12M" in text


def test_parse_docx_without_vlm_skips_image_captions(tmp_path):
    path = tmp_path / "report.docx"
    _build_docx(path)

    with patch.object(parsers, "_caption_image") as mock_caption:
        parsers.parse_docx(path, vlm=None)
        # No VLM config — captioning helper must not be called even if the
        # DOCX had images (this one doesn't, but the assertion guards the
        # default-off path).
        mock_caption.assert_not_called()


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def test_parse_pptx_one_page_per_slide(tmp_path):
    path = tmp_path / "deck.pptx"
    _build_pptx(path)

    doc = parsers.parse_pptx(path)
    assert doc.file_type == "pptx"
    assert len(doc.pages) == 2
    assert doc.pages[0].content.startswith("--- Slide 1 ---")
    assert "Intro" in doc.pages[0].content
    assert "Welcome to the deck." in doc.pages[0].content
    assert "Numbers" in doc.pages[1].content
    assert "Revenue doubled YoY." in doc.pages[1].content


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------

def test_parse_xlsx_groups_by_sheet(tmp_path):
    path = tmp_path / "metrics.xlsx"
    _build_xlsx(path)

    doc = parsers.parse_xlsx(path)
    text = doc.as_text()
    assert "--- Sheet: Q1 ---" in text
    assert "Metric | Value" in text
    assert "ARR | $12M" in text
    assert "--- Sheet: Q2 ---" in text


# ---------------------------------------------------------------------------
# Dispatch + fallback
# ---------------------------------------------------------------------------

def test_extract_document_text_dispatches_by_suffix(tmp_path):
    docx_path = tmp_path / "a.docx"
    _build_docx(docx_path)
    pptx_path = tmp_path / "b.pptx"
    _build_pptx(pptx_path)

    assert "# Quarterly Report" in parsers.extract_document_text(docx_path)
    assert "--- Slide 1 ---" in parsers.extract_document_text(pptx_path)


def test_extract_document_text_falls_back_to_plain_text(tmp_path):
    path = tmp_path / "note.txt"
    path.write_text("hello world", encoding="utf-8")
    assert parsers.extract_document_text(path) == "hello world"


def test_can_parse_recognises_office_suffixes():
    assert parsers.can_parse(".pdf")
    assert parsers.can_parse(".docx")
    assert parsers.can_parse(".doc")
    assert parsers.can_parse(".pptx")
    assert parsers.can_parse(".ppt")
    assert parsers.can_parse(".xlsx")
    assert parsers.can_parse(".xls")
    assert not parsers.can_parse(".txt")
    assert not parsers.can_parse(".md")


# ---------------------------------------------------------------------------
# Legacy format conversion
# ---------------------------------------------------------------------------

def test_legacy_doc_raises_helpful_error_when_soffice_missing(tmp_path):
    legacy = tmp_path / "old.doc"
    legacy.write_bytes(b"\xd0\xcf\x11\xe0fake-ole-header")  # not a real .doc

    with patch.object(parsers, "find_soffice", return_value=None):
        with pytest.raises(parsers.LegacyOfficeConversionError) as exc_info:
            parsers.extract_document_text(legacy)
        # Error message must point the user at LibreOffice — that's the
        # whole reason this exception is surfaced rather than swallowed.
        assert "LibreOffice" in str(exc_info.value)


def test_legacy_doc_runs_soffice_when_available(tmp_path):
    """When soffice is on PATH the converter should be invoked. We mock the
    subprocess call so the test doesn't require LibreOffice to be installed."""
    legacy = tmp_path / "old.doc"
    legacy.write_bytes(b"fake")

    converted_path_holder: dict[str, Path] = {}

    def fake_run(cmd, capture_output, text, timeout):  # noqa: ARG001
        out_dir = Path(cmd[cmd.index("--outdir") + 1])
        out_dir.mkdir(parents=True, exist_ok=True)
        produced = out_dir / f"{legacy.stem}.docx"
        # Drop a minimal valid DOCX into the output directory.
        _build_docx(produced)
        converted_path_holder["path"] = produced

        class Result:
            returncode = 0
            stderr = ""
            stdout = ""

        return Result()

    with (
        patch.object(parsers, "find_soffice", return_value="/usr/bin/soffice"),
        patch.object(parsers.subprocess, "run", side_effect=fake_run),
    ):
        text = parsers.extract_document_text(legacy)

    assert "# Quarterly Report" in text
    # Conversion artifact + its parent tempdir get cleaned up after parsing.
    assert not converted_path_holder["path"].exists()


# ---------------------------------------------------------------------------
# VLM integration smoke (mocked)
# ---------------------------------------------------------------------------

def test_pdf_page_complex_uses_fitz_bbox_heuristic():
    """The complexity helper takes a fitz Page object now. We fake it with a
    SimpleNamespace exposing the three attributes the helper touches:
    ``get_text``, ``get_images``, and ``rect``."""
    from types import SimpleNamespace

    class FakeRect:
        def __init__(self, w: float, h: float):
            self.width = w
            self.height = h

    def make_page(text: str, images: list, *, rect=FakeRect(600.0, 800.0), bbox_for=None):
        bbox_for = bbox_for or {}

        def get_image_bbox(img):
            if img in bbox_for:
                return bbox_for[img]
            return FakeRect(10, 10)

        return SimpleNamespace(
            get_text=lambda: text,
            get_images=lambda full=False: images,
            rect=rect,
            get_image_bbox=get_image_bbox,
        )

    # No images → never complex regardless of text
    assert parsers._pdf_page_is_complex(make_page("", [])) is False
    assert parsers._pdf_page_is_complex(make_page("text " * 50, [])) is False

    # No text + at least one image → complex
    assert parsers._pdf_page_is_complex(make_page("", [("xref-1",)])) is True

    # Plenty of text → not complex even with images
    assert (
        parsers._pdf_page_is_complex(
            make_page("a" * 1000, [("xref-1",)])
        )
        is False
    )

    # Thin text + a big image (covers > 5% of the page) → complex
    big_bbox = FakeRect(200, 200)  # 40000 / 480000 = ~8.3%
    img = ("xref-2",)
    assert (
        parsers._pdf_page_is_complex(
            make_page("short", [img], bbox_for={img: big_bbox})
        )
        is True
    )

    # Thin text + only a tiny image (<5%) → not complex
    tiny_bbox = FakeRect(20, 20)  # 400 / 480000 = ~0.08%
    img2 = ("xref-3",)
    assert (
        parsers._pdf_page_is_complex(
            make_page("short", [img2], bbox_for={img2: tiny_bbox})
        )
        is False
    )


# ---------------------------------------------------------------------------
# Parallel VLM captioning + integer/float xlsx formatting
# ---------------------------------------------------------------------------


def test_caption_images_batch_runs_in_parallel(monkeypatch):
    """Ensure batch dispatch hits every task and returns results keyed by id.
    We patch ``_caption_image`` to a deterministic stub so concurrency is
    observable but not racy."""
    calls: list[str] = []

    def stub(image_bytes, vlm, *, context_hint, strict_filtering):  # noqa: ARG001
        calls.append(context_hint)
        return (f"caption-for-{context_hint}", 7)

    monkeypatch.setattr(parsers, "_caption_image", stub)

    vlm = parsers.VLMConfig(model="m", api_key="k", max_workers=4)
    tasks = [(f"id{i}", b"\x89PNG", f"ctx{i}", False) for i in range(6)]
    results = parsers._caption_images_batch(tasks, vlm)

    assert set(results.keys()) == {f"id{i}" for i in range(6)}
    for task_id, (caption, tokens) in results.items():
        assert caption == f"caption-for-ctx{task_id[2:]}"
        assert tokens == 7
    assert sorted(calls) == [f"ctx{i}" for i in range(6)]


def test_caption_images_batch_single_worker_falls_back_to_sequential():
    """When max_workers <= 1 we should not spin a thread pool."""
    vlm = parsers.VLMConfig(model="m", api_key="k", max_workers=1)
    with patch.object(parsers, "_caption_image", return_value=("ok", 3)) as mock:
        results = parsers._caption_images_batch(
            [("a", b"x", "h", False), ("b", b"y", "h", False)], vlm
        )
    assert results == {"a": ("ok", 3), "b": ("ok", 3)}
    assert mock.call_count == 2


def test_parse_xlsx_handles_merged_cells_dates_and_int_floats(tmp_path):
    """Verify the xlsx accuracy improvements: merged header forward-fill,
    datetime → ``YYYY-MM-DD`` for date-only values, and ``12.0`` → ``12``."""
    from datetime import datetime

    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "对账"
    ws.append(["2024 销售对比", None, None, "备注"])
    ws.merge_cells("A1:C1")
    ws.append(["品类", "Q1", "Q2", "上半年"])
    ws.append(["电子产品", 320.0, 410.0, 730])
    ws.append(["更新于", datetime(2024, 6, 30), None, None])
    path = tmp_path / "metrics.xlsx"
    wb.save(str(path))

    out = parsers.parse_xlsx(path).as_text()
    assert "2024 销售对比 | 2024 销售对比 | 2024 销售对比 | 备注" in out
    assert "电子产品 | 320 | 410 | 730" in out  # 320.0 → 320
    assert "更新于 | 2024-06-30" in out  # date-only no 00:00:00 tail
    # Trailing empty columns trimmed off the date row.
    assert "更新于 | 2024-06-30 | |" not in out
