"""Tests for presentations skill batch 2: chart / image / table slides."""
from __future__ import annotations

import importlib.util
import json
import struct
import subprocess
import sys
import zlib
from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402

SKILL_DIR = Path(__file__).resolve().parent.parent / "tokenmind" / "skills" / "presentations" / "scripts"


def _load_script(path: Path):
    spec = importlib.util.spec_from_file_location(path.stem, path)
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    sys.modules[path.stem] = module
    spec.loader.exec_module(module)
    return module


@pytest.fixture(scope="module")
def add_chart_slide():
    return _load_script(SKILL_DIR / "add_chart_slide.py")


@pytest.fixture(scope="module")
def add_image_slide():
    return _load_script(SKILL_DIR / "add_image_slide.py")


@pytest.fixture(scope="module")
def add_table_slide():
    return _load_script(SKILL_DIR / "add_table_slide.py")


def _make_deck_with_slide(tmp_path: Path) -> Path:
    """Deck with one Blank-layout slide ready to receive shapes."""
    out = tmp_path / "deck.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    prs.save(out)
    return out


def _make_tiny_png(path: Path) -> Path:
    """Write a minimal valid 1x1 PNG so add_picture has a real image to ingest."""
    sig = b"\x89PNG\r\n\x1a\n"

    def chunk(tag: bytes, data: bytes) -> bytes:
        crc = zlib.crc32(tag + data) & 0xFFFFFFFF
        return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", crc)

    ihdr = chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
    raw = b"\x00\xFF\x00\x00"  # filter=0, RGB pixel (255,0,0)
    idat = chunk(b"IDAT", zlib.compress(raw))
    iend = chunk(b"IEND", b"")
    path.write_bytes(sig + ihdr + idat + iend)
    return path


# --- add_chart_slide --------------------------------------------------------


def test_chart_bar_two_series(tmp_path, add_chart_slide):
    deck = _make_deck_with_slide(tmp_path)
    add_chart_slide.add_chart_slide(
        deck,
        slide_index=0,
        left_in=0.5, top_in=1.0, width_in=9, height_in=5,
        chart_type="bar",
        data={
            "categories": ["Q1", "Q2", "Q3", "Q4"],
            "series": {"Revenue": [120, 140, 170, 200], "Costs": [90, 100, 110, 125]},
        },
        title="FY24",
    )
    prs = Presentation(deck)
    slide = prs.slides[0]
    charts = [s for s in slide.shapes if s.has_chart]
    assert len(charts) == 1
    chart = charts[0].chart
    series_names = [s.name for s in chart.series]
    assert series_names == ["Revenue", "Costs"]
    assert chart.has_title
    assert chart.chart_title.text_frame.text == "FY24"


def test_chart_line_single_series(tmp_path, add_chart_slide):
    deck = _make_deck_with_slide(tmp_path)
    add_chart_slide.add_chart_slide(
        deck,
        slide_index=0,
        left_in=0.5, top_in=1.0, width_in=9, height_in=5,
        chart_type="line",
        data={"categories": ["2020", "2021", "2022"], "series": {"ARR": [12, 18, 27]}},
    )
    prs = Presentation(deck)
    charts = [s for s in prs.slides[0].shapes if s.has_chart]
    assert len(charts) == 1


def test_chart_pie_drops_extra_series(tmp_path, add_chart_slide):
    deck = _make_deck_with_slide(tmp_path)
    add_chart_slide.add_chart_slide(
        deck,
        slide_index=0,
        left_in=1, top_in=1, width_in=8, height_in=5,
        chart_type="pie",
        data={
            "categories": ["A", "B", "C"],
            "series": {"Mix": [55, 30, 15], "Ignored": [1, 2, 3]},
        },
    )
    prs = Presentation(deck)
    chart = next(s for s in prs.slides[0].shapes if s.has_chart).chart
    # Pie keeps only the first series
    assert [s.name for s in chart.series] == ["Mix"]


def test_chart_invalid_type(tmp_path, add_chart_slide):
    deck = _make_deck_with_slide(tmp_path)
    with pytest.raises(ValueError, match="--type"):
        add_chart_slide.add_chart_slide(
            deck, slide_index=0, left_in=0, top_in=0, width_in=1, height_in=1,
            chart_type="donut",
            data={"categories": ["A"], "series": {"x": [1]}},
        )


def test_chart_series_length_mismatch(tmp_path, add_chart_slide):
    deck = _make_deck_with_slide(tmp_path)
    with pytest.raises(ValueError, match="does not match"):
        add_chart_slide.add_chart_slide(
            deck, slide_index=0, left_in=0, top_in=0, width_in=1, height_in=1,
            chart_type="bar",
            data={"categories": ["A", "B", "C"], "series": {"x": [1, 2]}},
        )


def test_chart_missing_keys(tmp_path, add_chart_slide):
    deck = _make_deck_with_slide(tmp_path)
    with pytest.raises(ValueError, match="categories"):
        add_chart_slide.add_chart_slide(
            deck, slide_index=0, left_in=0, top_in=0, width_in=1, height_in=1,
            chart_type="bar",
            data={"series": {"x": [1, 2]}},
        )


def test_chart_empty_series(tmp_path, add_chart_slide):
    deck = _make_deck_with_slide(tmp_path)
    with pytest.raises(ValueError, match="non-empty"):
        add_chart_slide.add_chart_slide(
            deck, slide_index=0, left_in=0, top_in=0, width_in=1, height_in=1,
            chart_type="bar",
            data={"categories": ["A"], "series": {}},
        )


# --- add_image_slide --------------------------------------------------------


def test_image_basic_placement(tmp_path, add_image_slide):
    deck = _make_deck_with_slide(tmp_path)
    img = _make_tiny_png(tmp_path / "logo.png")
    add_image_slide.add_image_slide(
        deck, slide_index=0, image=img, left_in=0.4, top_in=0.4, height_in=1.0,
    )
    prs = Presentation(deck)
    pics = [s for s in prs.slides[0].shapes if s.shape_type == 13]  # PICTURE
    assert len(pics) == 1


def test_image_forced_dimensions(tmp_path, add_image_slide):
    deck = _make_deck_with_slide(tmp_path)
    img = _make_tiny_png(tmp_path / "hero.png")
    add_image_slide.add_image_slide(
        deck, slide_index=0, image=img,
        left_in=0, top_in=0, width_in=13.333, height_in=7.5,
    )
    prs = Presentation(deck)
    pic = next(s for s in prs.slides[0].shapes if s.shape_type == 13)
    from pptx.util import Inches
    assert pic.width == Inches(13.333)
    assert pic.height == Inches(7.5)


def test_image_missing_file(tmp_path, add_image_slide):
    deck = _make_deck_with_slide(tmp_path)
    with pytest.raises(FileNotFoundError):
        add_image_slide.add_image_slide(
            deck, slide_index=0, image=tmp_path / "nope.png",
            left_in=0, top_in=0,
        )


def test_image_invalid_slide(tmp_path, add_image_slide):
    deck = _make_deck_with_slide(tmp_path)
    img = _make_tiny_png(tmp_path / "x.png")
    with pytest.raises(ValueError, match="--slide"):
        add_image_slide.add_image_slide(
            deck, slide_index=99, image=img, left_in=0, top_in=0,
        )


# --- add_table_slide --------------------------------------------------------


def test_table_basic_with_header(tmp_path, add_table_slide):
    deck = _make_deck_with_slide(tmp_path)
    add_table_slide.add_table_slide(
        deck,
        slide_index=0,
        left_in=0.5, top_in=2, width_in=9, height_in=2,
        data=[
            ["Metric", "Q3", "Q4"],
            ["Revenue", "$12.1M", "$14.8M"],
            ["Margin", "42%", "47%"],
        ],
        header_row=True,
    )
    prs = Presentation(deck)
    tables = [s for s in prs.slides[0].shapes if s.has_table]
    assert len(tables) == 1
    table = tables[0].table
    assert len(table.rows) == 3
    assert len(table.columns) == 3
    # Header row was styled bold
    hdr_cell = table.cell(0, 0)
    bold_flags = [r.font.bold for p in hdr_cell.text_frame.paragraphs for r in p.runs]
    assert any(bold_flags)
    # Body row is NOT bold
    body_cell = table.cell(1, 0)
    body_bold_flags = [r.font.bold for p in body_cell.text_frame.paragraphs for r in p.runs]
    assert not any(body_bold_flags)


def test_table_header_colors(tmp_path, add_table_slide):
    deck = _make_deck_with_slide(tmp_path)
    add_table_slide.add_table_slide(
        deck,
        slide_index=0,
        left_in=0.5, top_in=2, width_in=9, height_in=2,
        data=[["A", "B"], ["1", "2"]],
        header_row=True,
        header_bg="1F2937",
        header_font_color="FFFFFF",
    )
    prs = Presentation(deck)
    table = next(s for s in prs.slides[0].shapes if s.has_table).table
    hdr_cell = table.cell(0, 0)
    fill_color = hdr_cell.fill.fore_color.rgb
    assert str(fill_color) == "1F2937"
    text_color = hdr_cell.text_frame.paragraphs[0].runs[0].font.color.rgb
    assert str(text_color) == "FFFFFF"


def test_table_jagged_rows_rejected(tmp_path, add_table_slide):
    deck = _make_deck_with_slide(tmp_path)
    with pytest.raises(ValueError, match="same column count"):
        add_table_slide.add_table_slide(
            deck, slide_index=0, left_in=0, top_in=0, width_in=1, height_in=1,
            data=[["A", "B"], ["1"]],
        )


def test_table_empty_data_rejected(tmp_path, add_table_slide):
    deck = _make_deck_with_slide(tmp_path)
    with pytest.raises(ValueError, match="non-empty"):
        add_table_slide.add_table_slide(
            deck, slide_index=0, left_in=0, top_in=0, width_in=1, height_in=1,
            data=[],
        )


def test_table_non_array_data_rejected(tmp_path, add_table_slide):
    deck = _make_deck_with_slide(tmp_path)
    with pytest.raises(ValueError, match="2D array"):
        add_table_slide.add_table_slide(
            deck, slide_index=0, left_in=0, top_in=0, width_in=1, height_in=1,
            data={"not": "a list"},  # type: ignore[arg-type]
        )


def test_table_handles_non_string_cells(tmp_path, add_table_slide):
    deck = _make_deck_with_slide(tmp_path)
    add_table_slide.add_table_slide(
        deck,
        slide_index=0,
        left_in=0, top_in=0, width_in=9, height_in=2,
        data=[["Name", "Qty"], ["Widget", 10], ["Gadget", 7.5]],
    )
    prs = Presentation(deck)
    table = next(s for s in prs.slides[0].shapes if s.has_table).table
    assert table.cell(1, 1).text == "10"
    assert table.cell(2, 1).text == "7.5"


def test_table_data_file(tmp_path, add_table_slide):
    deck = _make_deck_with_slide(tmp_path)
    data_file = tmp_path / "table.json"
    data_file.write_text(
        json.dumps([["a", "b"], ["1", "2"], ["3", "4"]]),
        encoding="utf-8",
    )
    # Tested via CLI (the function takes data directly)
    r = subprocess.run(
        [
            sys.executable, str(SKILL_DIR / "add_table_slide.py"), str(deck),
            "--slide", "0",
            "--left", "0", "--top", "0", "--width", "5", "--height", "2",
            "--data-file", str(data_file),
        ],
        capture_output=True, text=True,
    )
    assert r.returncode == 0, r.stderr
    prs = Presentation(deck)
    table = next(s for s in prs.slides[0].shapes if s.has_table).table
    assert len(table.rows) == 3 and len(table.columns) == 2


# --- end-to-end CLI smoke ---------------------------------------------------


def test_cli_smoke_chart_image_table(tmp_path):
    deck = _make_deck_with_slide(tmp_path)
    img = _make_tiny_png(tmp_path / "logo.png")

    # Chart
    r = subprocess.run(
        [
            sys.executable, str(SKILL_DIR / "add_chart_slide.py"), str(deck),
            "--slide", "0", "--type", "bar",
            "--left", "0.5", "--top", "1", "--width", "6", "--height", "4",
            "--data", json.dumps({
                "categories": ["A", "B", "C"],
                "series": {"X": [1, 2, 3]},
            }),
            "--title", "Demo",
        ],
        capture_output=True, text=True,
    )
    assert r.returncode == 0, r.stderr

    # Image
    r = subprocess.run(
        [
            sys.executable, str(SKILL_DIR / "add_image_slide.py"), str(deck),
            "--slide", "0", "--image", str(img),
            "--left", "10", "--top", "0.5", "--height", "1",
        ],
        capture_output=True, text=True,
    )
    assert r.returncode == 0, r.stderr

    # Table
    r = subprocess.run(
        [
            sys.executable, str(SKILL_DIR / "add_table_slide.py"), str(deck),
            "--slide", "0",
            "--left", "0.5", "--top", "5.5", "--width", "12", "--height", "1.5",
            "--data", json.dumps([["k", "v"], ["one", "1"], ["two", "2"]]),
            "--header-row",
        ],
        capture_output=True, text=True,
    )
    assert r.returncode == 0, r.stderr

    prs = Presentation(deck)
    slide = prs.slides[0]
    assert any(s.has_chart for s in slide.shapes)
    assert any(s.shape_type == 13 for s in slide.shapes)
    assert any(s.has_table for s in slide.shapes)
